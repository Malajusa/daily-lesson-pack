#!/usr/bin/env python3
"""Screen a DLP deck for the approved visual-exemplar grammar."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation


SKILL_ROOT = Path(__file__).resolve().parents[1]
EXEMPLAR = SKILL_ROOT / "assets" / "visual-exemplars" / "t3w6-tuesday-edited-visual-exemplar.pptx"
EXEMPLAR_SHA256 = "069ec730ef879bafd25f5aae3e8d3e9ca6378ace76a46dfe0fea4897c857e880"
ROLE_COLOURS = {
    "amber": {"FFF3BF", "D6A900"},
    "blue": {"EAF2F8", "005A9C"},
    "green": {"EAF7EE", "1B7F3A"},
}
REQUIRED_REVIEW_CHECKS = {
    "VISUAL.HIERARCHY",
    "VISUAL.PROJECTION_READABILITY",
    "VISUAL.SPACE_USE",
    "VISUAL.REPRESENTATION_CLARITY",
    "VISUAL.CROSS_SLIDE_CONSISTENCY",
}
GENERIC_EVIDENCE = re.compile(r"^(?:checked|reviewed|intentional|looks good|acceptable|pass)$", re.I)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def rgb(value) -> str | None:
    try:
        colour = value.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(colour) if colour else None


def inspect_deck(path: Path) -> dict:
    deck = Presentation(path)
    fills = Counter()
    lines = Counter()
    fonts = Counter()
    rail_slides = set()
    slides_with_panels = set()

    for slide_number, slide in enumerate(deck.slides, start=1):
        significant_panels = 0
        for shape in slide.shapes:
            if getattr(shape, "fill", None) is not None and shape.fill.type:
                try:
                    foreground = shape.fill.fore_color
                except (AttributeError, TypeError, ValueError):
                    foreground = None
                colour = rgb(foreground)
                if colour:
                    fills[colour] += 1
                if shape.width * shape.height >= deck.slide_width * deck.slide_height * 0.04:
                    significant_panels += 1
            if getattr(shape, "line", None) is not None:
                try:
                    line_colour = shape.line.color
                except (AttributeError, TypeError, ValueError):
                    line_colour = None
                colour = rgb(line_colour)
                if colour:
                    lines[colour] += 1
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
            if (
                shape.left <= deck.slide_width * 0.02
                and shape.width <= deck.slide_width * 0.03
                and shape.height >= deck.slide_height * 0.80
            ):
                rail_slides.add(slide_number)
        if significant_panels >= 1:
            slides_with_panels.add(slide_number)

    role_counts = {
        name: sum(fills[colour] + lines[colour] for colour in colours)
        for name, colours in ROLE_COLOURS.items()
    }
    total_font_runs = sum(fonts.values())
    trebuchet_share = fonts["Trebuchet MS"] / total_font_runs if total_font_runs else 0.0
    slide_count = len(deck.slides)
    return {
        "slide_count": slide_count,
        "rail_coverage": len(rail_slides) / slide_count if slide_count else 0.0,
        "panel_coverage": len(slides_with_panels) / slide_count if slide_count else 0.0,
        "role_colour_counts": role_counts,
        "dominant_fonts": fonts.most_common(8),
        "trebuchet_share": round(trebuchet_share, 4),
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", type=Path, required=True)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument(
        "--review-record",
        type=Path,
        help="Independent rendered-review JSON bound to this deck SHA-256",
    )
    args = parser.parse_args()

    failures = []
    if not EXEMPLAR.is_file():
        failures.append("Approved visual exemplar is missing.")
        exemplar_hash = None
    else:
        exemplar_hash = sha256(EXEMPLAR)
        if exemplar_hash != EXEMPLAR_SHA256:
            failures.append("Approved visual exemplar checksum does not match the registered file.")

    metrics = inspect_deck(args.deck)
    if metrics["slide_count"] == 0:
        failures.append("Deck contains no slides.")
    if metrics["rail_coverage"] < 0.95:
        failures.append("Full-height left role rail appears on fewer than 95% of slides.")
    if metrics["panel_coverage"] < 0.90:
        failures.append("Significant instructional panel coverage is below 90% of slides.")
    if metrics["trebuchet_share"] < 0.80:
        failures.append("Trebuchet MS is not the dominant deck typeface at the required 80% threshold.")
    for role, count in metrics["role_colour_counts"].items():
        if count < 4:
            failures.append(f"The {role} role-colour family is not used often enough to establish the exemplar grammar.")

    deck_hash = sha256(args.deck)
    review_summary = {"provided": bool(args.review_record), "valid": False, "reviewer": None, "run_id": None, "generator_run_id": None}
    if not args.review_record or not args.review_record.is_file():
        failures.append("Independent rendered visual-review record is missing.")
    else:
        try:
            review = json.loads(args.review_record.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            review = {}
            failures.append(f"Independent visual-review record is unreadable: {exc}")
        reviewer = str(review.get("reviewer", "")).strip()
        run_id = str(review.get("run_id", "")).strip()
        generator_run_id = str(review.get("generator_run_id", "")).strip()
        review_summary["reviewer"] = reviewer or None
        review_summary["run_id"] = run_id or None
        review_summary["generator_run_id"] = generator_run_id or None
        if str(review.get("artifact_sha256", "")).lower() != deck_hash:
            failures.append("Independent visual review is not bound to this deck SHA-256.")
        if not reviewer or not run_id or not generator_run_id:
            failures.append("Independent visual review needs a reviewer, run_id and generator_run_id.")
        elif run_id == generator_run_id:
            failures.append("Independent visual review must use a different run_id from generation.")
        checks = review.get("checks")
        found: set[str] = set()
        if not isinstance(checks, list):
            failures.append("Independent visual review needs evidence-bearing check objects.")
            checks = []
        for check in checks:
            if not isinstance(check, dict):
                continue
            check_id = str(check.get("id", "")).strip()
            found.add(check_id)
            evidence = str(check.get("evidence", "")).strip()
            if str(check.get("result", "")).upper() != "PASS" or len(evidence) < 20 or GENERIC_EVIDENCE.fullmatch(evidence):
                failures.append(f"Visual review check {check_id or 'MISSING'} is not substantiated.")
        for missing in sorted(REQUIRED_REVIEW_CHECKS - found):
            failures.append(f"Independent visual review is missing {missing}.")
        review_summary["valid"] = not any("Independent visual review" in item or "visual-review record" in item or "Visual review check" in item for item in failures)

    report = {
        "status": "PASS" if not failures else "FAIL",
        "deck": str(args.deck.resolve()),
        "artifact_sha256": deck_hash,
        "exemplar": str(EXEMPLAR),
        "exemplar_sha256": exemplar_hash,
        "expected_exemplar_sha256": EXEMPLAR_SHA256,
        "metrics": metrics,
        "independent_review": review_summary,
        "failures": failures,
        "note": "Automated exemplar metrics are screening only. PASS requires a separate evidence-bearing rendered review.",
    }
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, indent=2))
    return 0 if not failures else 1


if __name__ == "__main__":
    raise SystemExit(main())
