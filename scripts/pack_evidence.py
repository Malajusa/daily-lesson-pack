"""Pack-wide evidence validation. Missing or unrecognised evidence fails closed.

Execution receipts are auditable provenance, not cryptographic proof of reviewer
independence. Capture them from the host; never author a fictitious receipt.
"""
from __future__ import annotations
import hashlib
import json
import re
import unicodedata
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]

def digest(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()

def read(path):
    return json.loads(Path(path).read_text())

def normal(text):
    return ' '.join(unicodedata.normalize('NFKC', text).split())

def shapes(items):
    for shape in items:
        if hasattr(shape, 'shapes'):
            yield from shapes(shape.shapes)
        else:
            yield shape

def pages(path):
    if path.suffix.lower() == '.pptx':
        return len(Presentation(path).slides)
    if path.suffix.lower() == '.pdf':
        import fitz
        with fitz.open(path) as doc:
            return len(doc)
    raise ValueError('Only PPTX/PDF deliverables are supported')

def extract(path, location):
    page = location['page']
    if not isinstance(page, int) or isinstance(page, bool) or page < 1:
        raise ValueError('Page must be a positive integer')
    if path.suffix.lower() == '.pptx':
        slide = Presentation(path).slides[page-1]
        matches = [s for s in shapes(slide.shapes) if s.shape_id == location['shape_id']]
        if len(matches) != 1 or not matches[0].has_text_frame:
            raise ValueError('Missing unique text shape')
        return matches[0].text
    import fitz
    box = location['bbox']
    if len(box) != 4 or box[0] >= box[2] or box[1] >= box[3]:
        raise ValueError('Invalid PDF text region')
    with fitz.open(path) as doc:
        return doc[page-1].get_text(clip=fitz.Rect(box))

def expected_checks(content, manifest, method):
    registry = read(ROOT/'references/qa-requirements.json')['requirements']
    expected = set()
    for rule in registry:
        if rule['method'] != method:
            continue
        if rule['scope'] == 'pack':
            expected.add((rule['id'], 'pack'))
        elif rule['scope'] == 'instance':
            expected.update((rule['id'], i['id']) for i in content['instances'] if i['owner'] == rule['owner'])
        elif rule['scope'] == 'task':
            expected.update((rule['id'], t['id']) for t in content['tasks'])
        elif rule['scope'] == 'page':
            expected.update((rule['id'], f"{a['id']}:{n}") for a in manifest['artifacts'] for n in range(1, a['pages']+1))
    return expected

def validate_content(content):
    errors = []
    if content.get('schema_version') != 3:
        errors.append('Content schema must be 3')
    instances = content['instances']
    tasks = content['tasks']
    if not instances or not tasks:
        errors.append('Content requires instances and tasks')
    ids = [i['id'] for i in instances]
    task_ids = [t['id'] for t in tasks]
    if len(ids) != len(set(ids)) or len(task_ids) != len(set(task_ids)):
        errors.append('Duplicate instance/task ID')
    owners = {i['id']: i['owner'] for i in instances}
    if set(owners) != {t['instance_id'] for t in tasks}:
        errors.append('Every scheduled content instance needs canonical tasks')
    for task in tasks:
        prefix = task['id']
        if task['instance_id'] not in owners:
            errors.append(prefix+': unknown instance')
        fields = task['fields']
        if not all(isinstance(v, str) and v.strip() for v in fields.values()):
            errors.append(prefix+': fields must be non-empty strings')
        if not fields.get('prompt') or not fields.get('answer'):
            errors.append(prefix+': prompt and answer required')
        demands = task.get('demands', [])
        if not demands or len({d['id'] for d in demands}) != len(demands):
            errors.append(prefix+': unique response demands required')
        for demand in demands:
            if not demand.get('action') or not demand.get('answer_quote') or normal(demand['answer_quote']) not in normal(fields['answer']):
                errors.append(prefix+': response demand lacks answer evidence')
        warmup = owners.get(task['instance_id']) == 'dlp-literacy-warmup'
        if warmup and len(demands) != 1 and not task.get('response_override_source'):
            errors.append(prefix+': Literacy warm-up requires one response')
        if warmup and re.search(r'\b(?:and\s+(?:explain|justify|name|write)|give\s+a\s+reason)\b', fields['prompt'], re.I) and not task.get('response_override_source'):
            errors.append(prefix+': additional response in prompt')
        if task.get('operation') == 'combine' and not re.search(r'\b(combine|join)\b', fields['prompt'], re.I):
            errors.append(prefix+': combination instruction lacks exact action')
        options = task.get('options', [])
        if options:
            if len(options) < 2 or sum(o.get('correct') is True for o in options) != 1:
                errors.append(prefix+': choice requires one correct option')
            for option in options:
                if not option.get('text') or normal(option['text']) not in normal(fields['prompt']):
                    errors.append(prefix+': option not present in prompt')
                if not option.get('correct') and (not option.get('misconception') or not option.get('rejection_reason')):
                    errors.append(prefix+': distractor lacks misconception evidence')
                if option.get('text', '').strip().lower() == 'for no reason':
                    errors.append(prefix+': known implausible distractor')
        if task.get('operation') == 'revise':
            before, after = fields.get('before', ''), fields.get('after', '')
            propositions = task.get('propositions', [])
            if not before or not after or not propositions:
                errors.append(prefix+': revision needs original, revision and proposition map')
            for fact in propositions:
                if not fact.get('before_quote') or not fact.get('after_quote') or normal(fact['before_quote']) not in normal(before) or normal(fact['after_quote']) not in normal(after):
                    errors.append(prefix+': proposition lost or unlocated')
    return errors

def audit_pack(manifest_path, content_path, context_path):
    manifest_path, content_path, context_path = map(Path, (manifest_path, content_path, context_path))
    manifest, content, context = map(read, (manifest_path, content_path, context_path))
    errors = validate_content(content)
    if manifest.get('schema_version') != 3:
        errors.append('Manifest schema must be 3')
    if manifest.get('content_sha256') != digest(content_path) or manifest.get('context_sha256') != digest(context_path):
        errors.append('Stale content/context binding')
    if content['instances'] != context.get('timetable_instances'):
        errors.append('Content instances differ from authoritative context')
    required = set(context.get('required_artifacts', []))
    if not {'deck','briefing'} <= required:
        errors.append('Context must explicitly declare required deliverables including deck and briefing')
    artifacts = manifest['artifacts']
    ids = [a['id'] for a in artifacts]
    roles = [a['role'] for a in artifacts]
    if len(ids) != len(set(ids)) or len(roles) != len(set(roles)) or set(roles) != required:
        errors.append('Missing, extra or duplicate deliverable')
    paths = {}
    for artifact in artifacts:
        path = (manifest_path.parent/artifact['path']).resolve()
        paths[artifact['id']] = path
        if digest(path) != artifact['sha256'] or pages(path) != artifact['pages']:
            errors.append(artifact['id']+': stale hash or page count')
    tasks = {t['id']:t for t in content['tasks']}
    docs = content.get('documents', {})
    seen, covered = set(), set()
    for binding in manifest['bindings']:
        key = (binding['artifact'], binding['record'], binding['field'])
        if key in seen:
            errors.append('Duplicate binding: '+str(key))
        seen.add(key)
        record = tasks.get(binding['record'], docs.get(binding['record']))
        expected = record['fields'][binding['field']]
        actual = extract(paths[binding['artifact']], binding)
        if normal(expected) != normal(actual):
            errors.append('Content mismatch: '+str(key))
        covered.add((binding['artifact'], binding['page']))
    for artifact in artifacts:
        for n in range(1,artifact['pages']+1):
            if (artifact['id'], n) not in covered:
                errors.append(f"Unbound page {artifact['id']}:{n}")
    # Every delivered text element must come from the shared source, including
    # titles/footers (store those in documents). One bound item per page is not
    # sufficient: otherwise unreviewed tasks could be added beside it.
    for artifact in artifacts:
        path = paths[artifact['id']]
        bindings = [b for b in manifest['bindings'] if b['artifact'] == artifact['id']]
        if path.suffix.lower() == '.pptx':
            bound_shapes = {(b['page'], b['shape_id']) for b in bindings}
            for n, slide in enumerate(Presentation(path).slides, 1):
                for shape in shapes(slide.shapes):
                    if getattr(shape, 'has_text_frame', False) and shape.text.strip() and (n,shape.shape_id) not in bound_shapes:
                        errors.append(f"Unbound slide text {artifact['id']}:{n}:{shape.shape_id}")
        else:
            import fitz
            with fitz.open(path) as doc:
                for n, page in enumerate(doc, 1):
                    regions = [fitz.Rect(b['bbox']) for b in bindings if b['page'] == n]
                    for word in page.get_text('words'):
                        centre = fitz.Point((word[0]+word[2])/2,(word[1]+word[3])/2)
                        if not any(centre in region for region in regions):
                            errors.append(f"Unbound PDF text {artifact['id']}:{n}: {word[4]}")
    by_role = {a['role']:a['id'] for a in artifacts}
    for task in content['tasks']:
        required_fields = [('deck','prompt'),('deck','answer')]
        if 'student' in required: required_fields.append(('student','prompt'))
        if 'answers' in required: required_fields.append(('answers','answer'))
        if task.get('operation') == 'revise':
            required_fields += [('deck','before'),('deck','after')]
        for role, field in required_fields:
            if (by_role.get(role),task['id'],field) not in seen:
                errors.append(f"Missing {role} {field}: {task['id']}")
    return errors, content, manifest

def audit_review(review_path, method, manifest_path, content_path, content, manifest, trace_path):
    review_path, trace_path = Path(review_path), Path(trace_path)
    review, trace = read(review_path), read(trace_path)
    errors = []
    if review.get('manifest_sha256') != digest(manifest_path) or review.get('content_sha256') != digest(content_path):
        errors.append(method+': review belongs to different pack/content')
    if review.get('requirements_sha256') != digest(ROOT/'references/qa-requirements.json'):
        errors.append(method+': review uses stale requirements')
    # The host exports this receipt. The generator must not fabricate its fields.
    if trace.get('source') not in ('collaboration','external-runner','human-review'):
        errors.append(method+': execution receipt source missing')
    if not trace.get('execution_id') or not trace.get('reviewer_actor') or not trace.get('generator_actor') or trace['reviewer_actor'] == trace['generator_actor']:
        errors.append(method+': separate reviewer execution missing')
    if trace.get('review_sha256') != digest(review_path) or trace.get('manifest_sha256') != digest(manifest_path):
        errors.append(method+': stale execution receipt')
    transcript = (trace_path.parent/trace['transcript_path']).resolve()
    if digest(transcript) != trace.get('transcript_sha256'):
        errors.append(method+': stale host transcript')
    if review.get('execution_id') != trace.get('execution_id'):
        errors.append(method+': review execution mismatch')
    tasks = {t['id']: t for t in content['tasks']}
    instances = {i['id'] for i in content['instances']}
    expected = expected_checks(content, manifest, method)
    found = set()
    artifact_ids = {a['id']:a for a in manifest['artifacts']}
    renders = {(r['artifact'],r['page']):r for r in manifest.get('renders',[])}
    for check in review['checks']:
        key = (check['id'],check['subject'])
        if key in found or key not in expected:
            errors.append(method+': duplicate or unexpected check '+str(key))
        found.add(key)
        # N/A is a reasoned applicability decision, never an absent entry.
        if check.get('result') not in ('PASS','NA') or not check.get('observation') or not check.get('citations'):
            errors.append(method+': incomplete/failed check '+str(key))
        if check.get('result') == 'NA' and (not check.get('applicability_reason') or method == 'visual'):
            errors.append(method+': invalid applicability decision '+str(key))
        if check.get('result') == 'NA' and check['id'].startswith('TASK.'):
            if check['id'] != 'TASK.DISTRACTORS' or tasks.get(check['subject'],{}).get('options'):
                errors.append('Required task check cannot be waived: '+str(key))
        for citation in check.get('citations', []):
            if check['subject'] in tasks or check['subject'] in instances:
                allowed_records = {check['subject']} if check['subject'] in tasks else {t['id'] for t in content['tasks'] if t['instance_id'] == check['subject']}
                if not any(b['record'] in allowed_records and b['artifact'] == citation['artifact'] and b['page'] == citation['page'] and b.get('shape_id') == citation.get('shape_id') and b.get('bbox') == citation.get('bbox') for b in manifest['bindings']):
                    errors.append('Citation does not cover the specified task/instance: '+str(key))
            artifact = artifact_ids[citation['artifact']]
            path = Path(manifest_path).parent/artifact['path']
            if citation.get('quote') and normal(citation['quote']) not in normal(extract(path,citation)):
                errors.append(method+': citation does not match final artefact')
            if not citation.get('quote') and method != 'visual':
                errors.append(method+': semantic citation requires exact quote')
            if method == 'visual':
                render = renders[(citation['artifact'],citation['page'])]
                render_path = Path(manifest_path).parent/render['path']
                if render.get('artifact_sha256') != artifact['sha256'] or digest(render_path) != render['sha256']:
                    errors.append('Visual render binding invalid')
                if check['subject'] != f"{citation['artifact']}:{citation['page']}":
                    errors.append('Visual citation covers wrong page')
    errors.extend(method+': missing check '+str(k) for k in sorted(expected-found))
    return errors
