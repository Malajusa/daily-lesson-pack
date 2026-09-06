#!/usr/bin/env python3
"""Audit deterministic Daily Lesson Pack release contracts.

This audit checks structural requirements that can be verified from a PPTX and
the pre-assembly component acceptance record. It does not replace semantic or
rendered visual review.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation


SKILL_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PROFILE = SKILL_ROOT / "references" / "default-pack-profile.json"
CONTENT_COMPONENTS = {
    "dlp-morning-work",
    "dlp-literacy-warmup",
    "dlp-shared-reading",
    "dlp-guided-reading",
    "dlp-writing-lesson",
    "dlp-numeracy-warmup",
    "dlp-maths-lesson",
}

WHITEBOARD_PATTERNS = (
    re.compile(r"\bwrite (?:one |your )?(?:complete )?(?:response|answer) on your whiteboard\b", re.I),
    re.compile(r"\buse (?:your|a) (?:mini[- ]?)?whiteboard\b", re.I),
    re.compile(r"\bshow (?:your )?(?:response|answer) on your whiteboard\b", re.I),
)
AMBIGUOUS_COMMA_RULE = re.compile(
    r"\b(?:place|put|add) (?:a )?comma(?:s)? after each (?:listed )?item except the last(?: one)?\b",
    re.I,
)
IMPRECISE_GRAMMAR_TERM = re.compile(r"\blinking words?\b", re.I)
UNQUOTED_DEFINITION = re.compile(r"(?<![\"“'‘])\b[A-Za-z][A-Za-z-]*\s*=")
GENERIC_SENTENCE_TASK = re.compile(r"\bwrite (?:a|one) sentence\b", re.I)
EXACT_COMBINATION_ACTION = re.compile(r"\b(?:combine|join)\b", re.I)
TWO_SENTENCE_TEXT = re.compile(r"[A-Z][^.!?]{3,}[.!?]\s+[A-Z][^.!?]{3,}[.!?]")
PROHIBITED_LITERACY_REASONING = re.compile(
    r"(?:^|[.!?]\s+)(?:please\s+)?(?:explain\s+(?:why|how(?:\s+you\s+know)?)|"
    r"how\s+do\s+you\s+know|justify\b)|\band\s+explain\s+(?:why|how(?:\s+you\s+know)?)\b",
    re.I | re.M,
)
PUNCTUATION_CUE = re.compile(
    r"\bAdded (?:punctuation|commas?|full stop|question mark|exclamation mark|apostrophe|quotation marks?|inverted commas?)\b",
    re.I,
)
UNQUOTED_METALANGUAGE = (
    re.compile(
        r"\b(?:use|using)\s+(?:the\s+)?(?:because|and|but|or|so)\b",
        re.I,
    ),
    re.compile(
        r"\b(?:word|conjunction|pronoun|prefix|suffix|preposition)\s+(?:because|its|hydro|bio)\b",
        re.I,
    ),
)
REMINDER_PANEL_COLOURS = {"FFF3BF", "FFF2CC", "FFE699"}
LITERACY_HEADER = re.compile(
    r"LITERACY\s+WARM[- ]?UP\s+(\d+)\s+OF\s+(\d+).*?\b(REMINDER|QUESTION|ANSWER)\b",
    re.I | re.S,
)
SHARED_HEADER = re.compile(
    r"SHARED\s+READING\s+(\d+)\s+OF\s+(\d+).*?\b(QUESTION|ANSWER)\b",
    re.I | re.S,
)
NUMERACY_HEADER = re.compile(
    r"NUMERACY\s+WARM[- ]?UP\s+(\d+)\s+OF\s+(\d+).*?\b(QUESTION|ANSWER)\b",
    re.I | re.S,
)
WARMUP_HEADER = re.compile(r"\b(?:LITERACY|NUMERACY|MATHEMATICS|MATHS)\s+WARM[- ]?UP\b", re.I)
MATHS_HEADER = re.compile(r"\b(?:NUMERACY|MATHEMATICS|MATHS)\b", re.I)
SLASH_FRACTION = re.compile(r"(?<![\w/])(\d{1,3})\s*/\s*(\d{1,3})(?![\w/])")
PROHIBITED_CONTEXT_SOURCE = re.compile(
    r"\b(?:chat memory|saved (?:personal )?context|project memory|"
    r"another account(?:'s)? project|standing (?:teaching )?preference)\b",
    re.I,
)
REQUIRED_CONTEXT_FIELDS = (
    "active_year_profile",
    "pack_profile",
    "date",
    "term_week",
    "day",
    "timetable",
    "mathematics_focus",
    "english_focus",
    "lesson_status",
)
REQUIRED_CHECKS = {
    "dlp-maths-lesson": {
        "MATHS.PLANNING",
        "MATHS.REPRESENTATION.PURPOSE",
        "MATHS.REPRESENTATION.EXACTNESS",
        "MATHS.DIFFERENTIATION",
        "MATHS.MODEL_PRACTICE",
        "MATHS.EXIT",
        "VISUAL.READABILITY",
    },
    "dlp-numeracy-warmup": {
        "NUMERACY.SEQUENCE",
        "NUMERACY.THREE_TASKS",
        "NUMERACY.ANSWERS",
        "VISUAL.READABILITY",
    },
}
GENERIC_REQUIRED_CHECKS = {"COMPONENT.CONTENT", "COMPONENT.ANSWERS", "VISUAL.READABILITY"}
GENERIC_EVIDENCE = re.compile(r"^(?:checked|reviewed|intentional|looks good|all passed|pass)$", re.I)


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_profile(path: Path) -> dict:
    data = json.loads(path.read_text(encoding="utf-8"))
    if data.get("schema_version") != 1:
        raise ValueError("classroom profile schema_version must be 1")
    pairs = data.get("numeracy_warmup", {}).get("prompt_answer_pairs")
    sequences = data.get("literacy_warmup", {}).get("sequence_count")
    if not isinstance(pairs, int) or pairs < 1 or not isinstance(sequences, int) or sequences < 1:
        raise ValueError("classroom profile must define positive warm-up counts")
    return data


def valid_evidence(value) -> bool:
    text = str(value or "").strip()
    return len(text) >= 12 and not GENERIC_EVIDENCE.fullmatch(text)


def normalise(text: str) -> str:
    return " ".join(text.split())


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def slide_text(slide) -> str:
    return "\n".join(filter(None, (shape_text(shape) for shape in slide.shapes)))


def shape_fill_rgb(shape) -> str | None:
    try:
        rgb = shape.fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb).upper() if rgb is not None else None


def run_rgb(run) -> str | None:
    try:
        rgb = run.font.color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb).upper() if rgb is not None else None


def is_green(rgb: str | None) -> bool:
    if not rgb or len(rgb) != 6:
        return False
    try:
        red, green, blue = (int(rgb[index:index + 2], 16) for index in (0, 2, 4))
    except ValueError:
        return False
    return green >= 90 and green > red * 1.15 and green > blue * 1.15


def shape_contains(outer, inner, tolerance: int = 1000) -> bool:
    return (
        int(outer.left) <= int(inner.left) + tolerance
        and int(outer.top) <= int(inner.top) + tolerance
        and int(outer.left + outer.width) >= int(inner.left + inner.width) - tolerance
        and int(outer.top + outer.height) >= int(inner.top + inner.height) - tolerance
    )


def remember_rule_in_yellow_panel(slide) -> bool:
    remember_shapes = [
        shape for shape in slide.shapes
        if re.search(r"\bRemember\s*:", shape_text(shape), re.I)
    ]
    for text_shape in remember_shapes:
        if shape_fill_rgb(text_shape) in REMINDER_PANEL_COLOURS:
            return True
        if any(
            panel is not text_shape
            and shape_fill_rgb(panel) in REMINDER_PANEL_COLOURS
            and shape_contains(panel, text_shape)
            for panel in slide.shapes
        ):
            return True
    return False


def expected_punctuation_chars(text: str) -> set[str]:
    lower = text.lower()
    expected: set[str] = set()
    pairs = (
        ("full stop", "."),
        ("comma", ","),
        ("question mark", "?"),
        ("exclamation mark", "!"),
        ("colon", ":"),
        ("semicolon", ";"),
        ("apostrophe", "'’"),
        ("quotation mark", '"“”'),
        ("inverted comma", '"“”'),
    )
    for label, chars in pairs:
        if label in lower:
            expected.update(chars)
    return expected


def has_highlighted_punctuation(slide, expected: set[str] | None = None) -> bool:
    expected = expected or set(".,!?;:'\"’“”")
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            runs = list(paragraph.runs)
            surrounding_sizes = [
                run.font.size.pt for run in runs
                if run.font.size is not None and any(char.isalnum() for char in run.text)
            ]
            if not surrounding_sizes and paragraph.font.size is not None:
                surrounding_sizes = [paragraph.font.size.pt]
            if not surrounding_sizes:
                continue
            base_size = min(surrounding_sizes)
            for run in runs:
                token = run.text.strip()
                if not token or not any(char in expected for char in token):
                    continue
                if any(char.isalnum() for char in token):
                    continue
                size = run.font.size.pt if run.font.size is not None else 0
                if run.font.bold is True and is_green(run_rgb(run)) and size >= base_size * 1.25 - 0.1:
                    return True
    return False


def issue(code: str, message: str, **details) -> dict:
    return {"severity": "fail", "code": code, "message": message, **details}


def audit_context_record(path: str | None) -> tuple[list[dict], dict]:
    issues: list[dict] = []
    summary = {"provided": bool(path), "resolved": 0, "required": len(REQUIRED_CONTEXT_FIELDS)}

    if not path:
        issues.append(issue(
            "context_record_missing",
            "No current-run context source record was supplied.",
        ))
        return issues, summary

    record_path = Path(path)
    if not record_path.is_file():
        issues.append(issue(
            "context_record_missing",
            "The supplied context source record does not exist.",
            path=str(record_path),
        ))
        return issues, summary

    try:
        data = json.loads(record_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        issues.append(issue(
            "context_record_invalid",
            "The context source record is not valid readable JSON.",
            path=str(record_path),
            detail=str(exc),
        ))
        return issues, summary

    if not isinstance(data, dict):
        issues.append(issue(
            "context_record_invalid",
            "The context source record must be a JSON object.",
        ))
        return issues, summary

    if data.get("schema_version") != 2:
        issues.append(issue("context_record_schema_outdated", "Context record must use schema version 2."))

    for field in REQUIRED_CONTEXT_FIELDS:
        entry = data.get(field)
        if not isinstance(entry, dict):
            issues.append(issue(
                "context_field_missing",
                "A required current-run context field is missing.",
                field=field,
            ))
            continue
        source = str(entry.get("source", "")).strip()
        resolved = entry.get("resolved") is True
        value = entry.get("value")
        if not resolved or value in (None, "", [], {}):
            issues.append(issue(
                "context_field_unresolved",
                "A required current-run context field is unresolved.",
                field=field,
            ))
        else:
            summary["resolved"] += 1
        if not source:
            issues.append(issue(
                "context_source_missing",
                "A required context field does not identify its current source.",
                field=field,
            ))
        elif PROHIBITED_CONTEXT_SOURCE.search(source):
            issues.append(issue(
                "prohibited_memory_source",
                "Chat memory or an unstated standing preference cannot authorise a pack field.",
                field=field,
                source=source,
            ))

    printing_requested = data.get("printing_requested") is True
    if printing_requested:
        printing = data.get("printing_quantity")
        if not isinstance(printing, dict) or printing.get("resolved") is not True:
            issues.append(issue(
                "printing_quantity_unresolved",
                "Printing was requested but the quantity is unresolved.",
            ))
        elif PROHIBITED_CONTEXT_SOURCE.search(str(printing.get("source", ""))):
            issues.append(issue(
                "prohibited_memory_source",
                "Chat memory cannot supply the printing quantity.",
                field="printing_quantity",
                source=str(printing.get("source", "")),
            ))

    day_entry = data.get("day")
    if isinstance(day_entry, dict) and "wednesday" in str(day_entry.get("value", "")).lower():
        for field in ("timetable", "mathematics_focus", "english_focus"):
            entry = data.get(field, {})
            if not isinstance(entry, dict) or entry.get("resolved") is not True:
                issues.append(issue(
                    "wednesday_bridge_unresolved",
                    "Wednesday requires a supplied timetable and day-level Mathematics and English focus.",
                    field=field,
                ))

    instances = data.get("timetable_instances")
    if not isinstance(instances, list) or not instances:
        issues.append(issue(
            "timetable_instances_missing",
            "Schema v2 requires concrete timetable instances with unique IDs and durations.",
        ))
    else:
        seen_ids: set[str] = set()
        for index, instance in enumerate(instances):
            if not isinstance(instance, dict):
                issues.append(issue("timetable_instance_invalid", "Each timetable instance must be an object.", instance_index=index))
                continue
            instance_id = str(instance.get("id", "")).strip()
            owner = str(instance.get("owner", "")).strip()
            if not instance_id or instance_id in seen_ids:
                issues.append(issue("timetable_instance_id_invalid", "Every timetable instance needs a unique non-empty ID.", instance_id=instance_id or "MISSING"))
            seen_ids.add(instance_id)
            if owner not in CONTENT_COMPONENTS:
                issues.append(issue("timetable_instance_owner_invalid", "A timetable instance names an unknown content owner.", instance_id=instance_id, owner=owner))
            duration = instance.get("duration_minutes")
            if not isinstance(duration, (int, float)) or duration <= 0:
                issues.append(issue("timetable_duration_invalid", "Every timetable instance needs a positive duration_minutes budget.", instance_id=instance_id))
            if not str(instance.get("start", "")).strip() or not str(instance.get("purpose", "")).strip():
                issues.append(issue("timetable_instance_incomplete", "Every timetable instance needs a start time and instructional purpose.", instance_id=instance_id))
        summary["timetable_instances"] = len(instances)

    return issues, summary


def audit_component_record(
    path: str | None,
    context_path: str | None = None,
    deck_path: str | None = None,
) -> tuple[list[dict], dict]:
    issues: list[dict] = []
    summary = {"provided": bool(path), "scheduled": 0, "recorded": 0, "passed": 0}

    if not path:
        issues.append(issue(
            "component_record_missing",
            "No pre-assembly component acceptance record was supplied.",
        ))
        return issues, summary

    record_path = Path(path)
    if not record_path.is_file():
        issues.append(issue(
            "component_record_missing",
            "The supplied component acceptance record does not exist.",
            path=str(record_path),
        ))
        return issues, summary

    try:
        data = json.loads(record_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        issues.append(issue(
            "component_record_invalid",
            "The component acceptance record is not valid readable JSON.",
            path=str(record_path),
            detail=str(exc),
        ))
        return issues, summary

    if not isinstance(data, dict):
        issues.append(issue(
            "component_record_invalid",
            "The component acceptance record must be a JSON object.",
            path=str(record_path),
        ))
        return issues, summary

    if data.get("schema_version") != 2:
        issues.append(issue(
            "component_record_schema_outdated",
            "Component acceptance must use schema version 2 with instance IDs and evidence-bearing checks.",
        ))
    generation_run_id = str(data.get("generation_run_id", "")).strip()
    summary["generation_run_id"] = generation_run_id or None
    if not generation_run_id:
        issues.append(issue("generation_run_id_missing", "Component evidence must identify the generation run."))

    scheduled = data.get("scheduled_instances")
    records = data.get("components")
    if not isinstance(scheduled, list) or not scheduled:
        issues.append(issue(
            "scheduled_instances_missing",
            "The acceptance record must list all scheduled component instances.",
        ))
        scheduled = []
    if not isinstance(records, list):
        issues.append(issue(
            "component_results_missing",
            "The acceptance record must contain a components result list.",
        ))
        records = []

    scheduled_objects = [item for item in scheduled if isinstance(item, dict)]
    if len(scheduled_objects) != len(scheduled):
        issues.append(issue("scheduled_instance_invalid", "Every scheduled instance must be a JSON object."))
    scheduled = scheduled_objects
    summary["scheduled"] = len(scheduled)
    summary["recorded"] = len(records)

    scheduled_ids = [str(item.get("id", "")).strip() for item in scheduled]
    for instance_id, count in Counter(scheduled_ids).items():
        if not instance_id or count > 1:
            issues.append(issue(
                "duplicate_scheduled_instance",
                "Scheduled instance IDs must be unique and non-empty.",
                instance_id=instance_id or "MISSING",
            ))
    scheduled_by_id = {str(item.get("id", "")).strip(): item for item in scheduled if str(item.get("id", "")).strip()}
    maths_instance_count = sum(str(item.get("owner", "")).strip() == "dlp-maths-lesson" for item in scheduled)
    fraction_focus = False
    active_year_profile = ""
    if context_path and Path(context_path).is_file():
        try:
            context_preview = json.loads(Path(context_path).read_text(encoding="utf-8"))
            focus = str(context_preview.get("mathematics_focus", {}).get("value", "")).lower()
            fraction_focus = "fraction" in focus and any(token in focus for token in ("decimal", "equivalent", "tenths", "hundredths"))
            active_year_profile = str(
                context_preview.get("active_year_profile", {}).get("value", "")
            ).strip()
        except (OSError, json.JSONDecodeError):
            pass
    for instance_id, item in scheduled_by_id.items():
        owner = str(item.get("owner", "")).strip()
        if owner not in CONTENT_COMPONENTS:
            issues.append(issue(
                "unknown_scheduled_component_owner",
                "The acceptance record names an unknown component owner.",
                instance_id=instance_id,
                owner=owner,
            ))

    by_id: dict[str, list[dict]] = defaultdict(list)
    for index, record in enumerate(records):
        if not isinstance(record, dict):
            issues.append(issue(
                "component_result_invalid",
                "Each component result must be a JSON object.",
                record_index=index,
            ))
            continue
        instance_id = str(record.get("instance_id", "")).strip()
        by_id[instance_id].append(record)

    for instance_id, matching in by_id.items():
        if not instance_id or instance_id not in scheduled_by_id:
            issues.append(issue(
                "unknown_component_result",
                "The acceptance record contains a result for an unknown instance.",
                instance_id=instance_id or "MISSING",
            ))
        if len(matching) > 1:
            issues.append(issue(
                "duplicate_component_result",
                "A component instance has multiple acceptance results.",
                instance_id=instance_id or "MISSING",
            ))

    for instance_id, scheduled_item in scheduled_by_id.items():
        owner = str(scheduled_item.get("owner", "")).strip()
        matching = by_id.get(instance_id, [])
        if not matching:
            issues.append(issue(
                "scheduled_component_result_missing",
                "A scheduled component has no acceptance result.",
                instance_id=instance_id,
            ))
            continue
        if len(matching) > 1:
            issues.append(issue(
                "duplicate_component_result",
                "A scheduled component has multiple acceptance results.",
                instance_id=instance_id,
            ))
            continue

        record = matching[0]
        record_owner = str(record.get("owner", "")).strip()
        if record_owner != owner:
            issues.append(issue("component_owner_mismatch", "The result owner does not match its scheduled instance.", instance_id=instance_id, expected=owner, actual=record_owner))
        status = str(record.get("status", "")).upper()
        checks = record.get("checks")
        artefact = record.get("artefact") or record.get("artifact") or record.get("slide_range")
        if status != "PASS":
            issues.append(issue(
                "component_not_passed",
                "Every scheduled component must record PASS before assembly.",
                instance_id=instance_id,
                status=status or "MISSING",
            ))
        else:
            summary["passed"] += 1
        if not isinstance(checks, list) or not checks:
            issues.append(issue(
                "component_checks_missing",
                "A component PASS must include evidence-bearing check objects.",
                instance_id=instance_id,
            ))
            checks = []
        check_ids: set[str] = set()
        for check in checks:
            if not isinstance(check, dict):
                issues.append(issue("component_check_invalid", "Checks must be objects, not freeform assertions.", instance_id=instance_id))
                continue
            check_id = str(check.get("id", "")).strip()
            check_ids.add(check_id)
            if str(check.get("result", "")).upper() != "PASS" or not valid_evidence(check.get("evidence")):
                issues.append(issue("component_check_unsubstantiated", "Every check needs a stable ID, PASS result and concrete evidence.", instance_id=instance_id, check_id=check_id or "MISSING"))
        required = set(REQUIRED_CHECKS.get(owner, GENERIC_REQUIRED_CHECKS))
        if owner == "dlp-maths-lesson" and active_year_profile == "year-4-5":
            required.update({"MATHS.YEAR4.PATHWAY", "MATHS.YEAR5.PATHWAY"})
        if owner == "dlp-maths-lesson" and maths_instance_count > 1:
            required.add("MATHS.BLOCK.BREAKPOINT")
        if owner == "dlp-maths-lesson" and fraction_focus:
            required.add("MATHS.FRACTION.REPARTITIONING")
        for missing in sorted(required - check_ids):
            issues.append(issue("component_required_check_missing", "A required component check is missing.", instance_id=instance_id, check_id=missing))
        if not artefact:
            issues.append(issue(
                "component_artefact_missing",
                "A component PASS must identify its artefact or slide range.",
                instance_id=instance_id,
            ))

        estimate = record.get("estimated_minutes")
        available = scheduled_item.get("duration_minutes")
        if not isinstance(estimate, (int, float)) or estimate <= 0:
            issues.append(issue("component_time_estimate_missing", "Every instance needs a positive estimated_minutes value.", instance_id=instance_id))
        elif isinstance(available, (int, float)) and estimate > available:
            issues.append(issue("component_time_budget_exceeded", "Estimated teaching time exceeds the scheduled duration.", instance_id=instance_id, estimated_minutes=estimate, available_minutes=available))

    if context_path and Path(context_path).is_file():
        try:
            context = json.loads(Path(context_path).read_text(encoding="utf-8"))
            context_instances = context.get("timetable_instances", [])
            expected = {(str(i.get("id", "")), str(i.get("owner", "")), i.get("duration_minutes")) for i in context_instances if isinstance(i, dict)}
            actual = {(str(i.get("id", "")), str(i.get("owner", "")), i.get("duration_minutes")) for i in scheduled if isinstance(i, dict)}
            if expected != actual:
                issues.append(issue("scheduled_instances_context_mismatch", "Scheduled instances do not exactly match the current context record.", expected=sorted(map(str, expected)), actual=sorted(map(str, actual))))
        except (OSError, json.JSONDecodeError):
            pass

    recorded_hash = str(data.get("artifact_sha256", "")).lower()
    if not recorded_hash:
        issues.append(issue("artifact_hash_missing", "The acceptance record must bind evidence to the assembled deck SHA-256."))
    elif deck_path and Path(deck_path).is_file() and recorded_hash != file_sha256(Path(deck_path)):
        issues.append(issue("artifact_hash_mismatch", "The acceptance evidence does not belong to the audited deck."))

    return issues, summary


def validate_sequence(
    entries: list[dict],
    roles: list[str],
    label: str,
    required_total: int | None = None,
) -> list[dict]:
    issues: list[dict] = []
    if not entries:
        if required_total is not None and required_total > 0:
            issues.append(issue(
                f"{label}_missing",
                f"{label.replace('_', ' ').title()} headers are missing.",
                expected_sequences=required_total,
            ))
        return issues

    totals = {entry["total"] for entry in entries}
    if len(totals) != 1:
        issues.append(issue(
            f"{label}_total_inconsistent",
            f"{label.replace('_', ' ').title()} headers use inconsistent sequence totals.",
            totals=sorted(totals),
        ))
    header_total = max(totals)
    expected_total = required_total if required_total is not None else header_total
    if required_total is not None and totals != {required_total}:
        issues.append(issue(
            f"{label}_count_invalid",
            f"{label.replace('_', ' ').title()} must contain the expected number of sequences.",
            expected=required_total,
            header_totals=sorted(totals),
        ))
    by_number: dict[int, list[dict]] = defaultdict(list)
    for entry in entries:
        by_number[entry["number"]].append(entry)

    for number in range(1, expected_total + 1):
        group = sorted(by_number.get(number, []), key=lambda item: item["slide"])
        if not group:
            issues.append(issue(
                f"{label}_sequence_missing",
                f"{label.replace('_', ' ').title()} sequence {number} is missing.",
                sequence=number,
            ))
            continue
        actual_roles = [entry["role"] for entry in group]
        actual_slides = [entry["slide"] for entry in group]
        if actual_roles != roles:
            issues.append(issue(
                f"{label}_role_order_invalid",
                f"{label.replace('_', ' ').title()} sequence {number} has the wrong role order.",
                sequence=number,
                expected=roles,
                actual=actual_roles,
                slides=actual_slides,
            ))
        if len(actual_slides) == len(roles) and any(
            later != earlier + 1 for earlier, later in zip(actual_slides, actual_slides[1:])
        ):
            issues.append(issue(
                f"{label}_slides_not_consecutive",
                f"{label.replace('_', ' ').title()} sequence {number} is not consecutive.",
                sequence=number,
                slides=actual_slides,
            ))

    expected_order = [
        (number, role)
        for number in range(1, expected_total + 1)
        for role in roles
    ]
    actual_order = [(entry["number"], entry["role"]) for entry in entries]
    if actual_order != expected_order:
        issues.append(issue(
            f"{label}_global_order_invalid",
            f"{label.replace('_', ' ').title()} sequences are not in complete numbered role order.",
            expected=expected_order,
            actual=actual_order,
        ))

    return issues


def audit_deck(path: str, literacy_count: int = 10, numeracy_count: int = 5, response_override_slides: set[int] | None = None) -> tuple[list[dict], dict]:
    try:
        presentation = Presentation(path)
    except Exception as exc:
        return [issue(
            "deck_unreadable",
            "The supplied deck could not be opened as a readable PPTX.",
            path=path,
            detail=str(exc),
        )], {
            "slides": 0,
            "literacy_sequence_slides": 0,
            "numeracy_sequence_slides": 0,
            "shared_reading_sequence_slides": 0,
        }
    issues: list[dict] = []
    literacy_entries: list[dict] = []
    numeracy_entries: list[dict] = []
    shared_entries: list[dict] = []

    for slide_number, slide in enumerate(presentation.slides, 1):
        full_text = slide_text(slide)
        is_warmup_slide = bool(WARMUP_HEADER.search(full_text))
        numeracy_match = NUMERACY_HEADER.search(full_text)
        if numeracy_match:
            numeracy_entries.append({
                "slide": slide_number,
                "number": int(numeracy_match.group(1)),
                "total": int(numeracy_match.group(2)),
                "role": numeracy_match.group(3).upper(),
            })
        if MATHS_HEADER.search(full_text):
            for match in SLASH_FRACTION.finditer(full_text):
                numerator, denominator = int(match.group(1)), int(match.group(2))
                if denominator and numerator <= denominator * 20:
                    issues.append(issue(
                        "projected_fraction_slash_notation",
                        "Projected Mathematics fractions must use stacked notation where practical.",
                        slide=slide_number,
                        text=match.group(0),
                    ))
        literacy_match = LITERACY_HEADER.search(full_text)
        if literacy_match:
            literacy_role = literacy_match.group(3).upper()
            literacy_entries.append({
                "slide": slide_number,
                "number": int(literacy_match.group(1)),
                "total": int(literacy_match.group(2)),
                "role": literacy_role,
            })
            if literacy_role == "REMINDER":
                if not re.search(r"\bRemember\s*:", full_text, re.I):
                    issues.append(issue(
                        "literacy_reminder_rule_missing",
                        "A Literacy reminder slide must contain a concise applicable Remember: rule.",
                        slide=slide_number,
                    ))
                elif not remember_rule_in_yellow_panel(slide):
                    issues.append(issue(
                        "literacy_reminder_rule_outside_yellow_panel",
                        "The Remember: rule is not contained by the yellow reminder panel.",
                        slide=slide_number,
                    ))
                expected = expected_punctuation_chars(full_text)
                if expected and not has_highlighted_punctuation(slide, expected):
                    issues.append(issue(
                        "literacy_reminder_punctuation_not_highlighted",
                        "A punctuation reminder must visibly enlarge, bold and colour the target mark in its example.",
                        slide=slide_number,
                    ))
            if literacy_role == "QUESTION" and re.search(r"\b(?:because|but|and|although|so|or|yet)\b", full_text, re.I):
                has_two_sentence_source = any(
                    TWO_SENTENCE_TEXT.search(shape_text(shape)) for shape in slide.shapes
                )
                if (
                    has_two_sentence_source
                    and GENERIC_SENTENCE_TASK.search(full_text)
                    and not EXACT_COMBINATION_ACTION.search(full_text)
                ):
                    issues.append(issue(
                        "literacy_combination_action_ambiguous",
                        "A sentence-combination task must explicitly say combine or join, not merely write a sentence.",
                        slide=slide_number,
                    ))
            if literacy_role == "QUESTION" and slide_number not in (response_override_slides or set()) and (PROHIBITED_LITERACY_REASONING.search(full_text) or re.search(r"\band\s+(?:name|write|give|complete)\b", full_text, re.I)):
                issues.append(issue(
                    "literacy_extra_reasoning_demand",
                    "A Literacy warm-up question must require one direct response without an added explain or justify demand.",
                    slide=slide_number,
                ))
            if literacy_role == "ANSWER" and PUNCTUATION_CUE.search(full_text):
                if not has_highlighted_punctuation(slide):
                    issues.append(issue(
                        "literacy_added_punctuation_not_highlighted",
                        "An answer that adds punctuation must contain a green, bold punctuation run at least 125% of the surrounding size.",
                        slide=slide_number,
                    ))
            if IMPRECISE_GRAMMAR_TERM.search(full_text):
                issues.append(issue(
                    "imprecise_grammar_term",
                    "Use the accurate grammatical term conjunction instead of linking word.",
                    slide=slide_number,
                ))
            definition_match = UNQUOTED_DEFINITION.search(full_text)
            if definition_match:
                issues.append(issue(
                    "unquoted_metalinguistic_word",
                    "A word being defined as language must be placed in inverted commas.",
                    slide=slide_number,
                    text=definition_match.group(0),
                ))
            for pattern in UNQUOTED_METALANGUAGE:
                match = pattern.search(full_text)
                if match:
                    issues.append(issue(
                        "unquoted_metalinguistic_word",
                        "A literal word discussed as language must be placed in inverted commas.",
                        slide=slide_number,
                        text=match.group(0),
                    ))
                    break
        shared_match = SHARED_HEADER.search(full_text)
        if shared_match:
            shared_entries.append({
                "slide": slide_number,
                "number": int(shared_match.group(1)),
                "total": int(shared_match.group(2)),
                "role": shared_match.group(3).upper(),
            })

        for shape in slide.shapes:
            text = normalise(shape_text(shape))
            if not text:
                continue
            lower_text = text.lower()
            is_footer_position = int(shape.top) >= int(presentation.slide_height * 0.70)
            if "whiteboard" in lower_text and (
                is_footer_position
                or (is_warmup_slide and any(pattern.search(text) for pattern in WHITEBOARD_PATTERNS))
            ):
                issues.append(issue(
                    "prohibited_whiteboard_footer",
                    "A generic whiteboard-response instruction appears on a warm-up/task slide.",
                    slide=slide_number,
                    shape=getattr(shape, "name", ""),
                    text=text[:240],
                ))
            if AMBIGUOUS_COMMA_RULE.search(text) and not re.search(
                r",\s+(?:and|or)\b", full_text, re.I
            ):
                issues.append(issue(
                    "ambiguous_comma_rule",
                    "The comma rule can contradict a model that omits the optional comma before and/or.",
                    slide=slide_number,
                    shape=getattr(shape, "name", ""),
                    text=text[:240],
                ))

    issues.extend(validate_sequence(
        literacy_entries,
        ["REMINDER", "QUESTION", "ANSWER"],
        "literacy_warmup",
        required_total=literacy_count,
    ))
    issues.extend(validate_sequence(
        numeracy_entries,
        ["QUESTION", "ANSWER"],
        "numeracy_warmup",
        required_total=numeracy_count,
    ))
    issues.extend(validate_sequence(
        shared_entries,
        ["QUESTION", "ANSWER"],
        "shared_reading",
    ))

    return issues, {
        "slides": len(presentation.slides),
        "literacy_sequence_slides": len(literacy_entries),
        "numeracy_sequence_slides": len(numeracy_entries),
        "shared_reading_sequence_slides": len(shared_entries),
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", required=True, help="PPTX Daily Lesson Pack to audit")
    parser.add_argument(
        "--component-record",
        help="JSON pre-assembly component acceptance record; absence is a failure",
    )
    parser.add_argument(
        "--context-record",
        help="JSON current-run context source record; absence is a failure",
    )
    parser.add_argument("--out", required=True, help="JSON report path")
    parser.add_argument(
        "--profile",
        type=Path,
        default=DEFAULT_PROFILE,
        help="Validated classroom profile JSON",
    )
    parser.add_argument(
        "--numeracy-count",
        type=int,
        default=None,
        help="Authorised Numeracy pair override; otherwise use the profile",
    )
    parser.add_argument(
        "--literacy-count",
        type=int,
        default=None,
        help="Authorised Literacy sequence override; otherwise use the profile",
    )
    parser.add_argument("--response-override-slides", type=int, nargs="*", default=[])
    args = parser.parse_args()

    try:
        profile = load_profile(args.profile)
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        parser.error(f"invalid classroom profile: {exc}")
    # Required counts follow resolved scheduled instances, never weekday guesses.
    try:
        runtime = json.loads(Path(args.context_record).read_text())
        owners = {i["owner"] for i in runtime["timetable_instances"]}
    except (OSError, TypeError, KeyError, ValueError):
        owners = {"dlp-literacy-warmup", "dlp-numeracy-warmup"}
    literacy_count = (args.literacy_count if args.literacy_count is not None else profile["literacy_warmup"]["sequence_count"]) if "dlp-literacy-warmup" in owners else 0
    numeracy_count = (args.numeracy_count if args.numeracy_count is not None else profile["numeracy_warmup"]["prompt_answer_pairs"]) if "dlp-numeracy-warmup" in owners else 0
    if ("dlp-literacy-warmup" in owners and literacy_count < 1) or ("dlp-numeracy-warmup" in owners and numeracy_count < 1):
        parser.error("scheduled sequence counts must be positive")

    context_issues, context_summary = audit_context_record(args.context_record)
    component_issues, component_summary = audit_component_record(
        args.component_record,
        context_path=args.context_record,
        deck_path=args.deck,
    )
    deck_issues, deck_summary = audit_deck(
        args.deck,
        literacy_count=literacy_count,
        numeracy_count=numeracy_count,
        response_override_slides=set(args.response_override_slides),
    )
    issues = context_issues + component_issues + deck_issues
    counts = Counter(item["code"] for item in issues)
    overall = "fail" if issues else "pass"
    report = {
        "deck": args.deck,
        "artifact_sha256": file_sha256(Path(args.deck)) if Path(args.deck).is_file() else None,
        "context_record": args.context_record,
        "component_record": args.component_record,
        "classroom_profile": str(args.profile),
        "overall_status": overall,
        "summary": {
            **deck_summary,
            "context_sources": context_summary,
            "component_acceptance": component_summary,
            "failures": len(issues),
            "failure_codes": dict(sorted(counts.items())),
        },
        "issues": issues,
        "note": "Deterministic contract audit only. Independent semantic and rendered visual QA remain mandatory.",
    }
    output = Path(args.out)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps({"overall_status": overall, **report["summary"]}))
    return 1 if issues else 0


if __name__ == "__main__":
    raise SystemExit(main())
