#!/usr/bin/env python3
"""Audit containment of student-facing text inside shaded/coloured PowerPoint panels.

The audit is intentionally conservative:
- hard failures are reported when a plausible panel/text pairing is strong and
  the text box crosses the panel boundary;
- inadequate margins on filled instructional text panels are hard failures;
- severe estimated text-fit overflow inside filled instructional panels is a hard failure;
- padding problems on explicitly owned panel/text pairs are hard failures;
- ambiguous spatial pairings are reported for human review.

Render-level inspection remains mandatory because PowerPoint line wrapping and
font metrics can create visible overflow even when nominal text-box bounds fit.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
DEFAULT_PADDING_IN = 0.15
MAX_PANEL_SLIDE_AREA = 0.92
MIN_PAIR_OVERLAP = 0.45
STRONG_PAIR_OVERLAP = 0.75
MIN_PANEL_HEAVY_SLIDES = 3
MIN_PANEL_HEAVY_COUNT = 5
PANEL_NAME_TOKENS = (
    "panel", "card", "footer", "box", "callout", "banner",
    "all_", "most_", "some_", "why_", "question_", "answer_",
)
TEXT_NAME_TOKENS = ("text", "task", "response", "prompt", "explanation", "label")


@dataclass
class Box:
    left: int
    top: int
    right: int
    bottom: int

    @property
    def width(self):
        return max(0, self.right - self.left)

    @property
    def height(self):
        return max(0, self.bottom - self.top)

    @property
    def area(self):
        return self.width * self.height

    @property
    def cx(self):
        return (self.left + self.right) / 2

    @property
    def cy(self):
        return (self.top + self.bottom) / 2


def iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def bbox(shape):
    return Box(
        int(shape.left),
        int(shape.top),
        int(shape.left + shape.width),
        int(shape.top + shape.height),
    )


def text_of(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def name_of(shape):
    return (getattr(shape, "name", "") or "").strip()


def has_visible_fill(shape):
    fill = getattr(shape, "fill", None)
    if fill is None:
        return False
    try:
        fill_type = fill.type
    except Exception:
        return False
    return fill_type is not None and fill_type != MSO_FILL.BACKGROUND


def is_candidate_panel(shape, slide_area):
    if shape.shape_type in {
        MSO_SHAPE_TYPE.LINE,
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.MEDIA,
    }:
        return False
    box = bbox(shape)
    if box.area <= 0 or box.area > slide_area * MAX_PANEL_SLIDE_AREA:
        return False
    if not has_visible_fill(shape):
        return False
    if box.width < 0.45 * EMU_PER_INCH or box.height < 0.28 * EMU_PER_INCH:
        return False
    return True


def is_text_candidate(shape):
    return getattr(shape, "has_text_frame", False) and len(text_of(shape).strip()) > 1


def intersection_area(a, b):
    left = max(a.left, b.left)
    top = max(a.top, b.top)
    right = min(a.right, b.right)
    bottom = min(a.bottom, b.bottom)
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def overlap_ratio(text_box, panel_box):
    return 0 if text_box.area == 0 else intersection_area(text_box, panel_box) / text_box.area


def centre_inside(text_box, panel_box):
    return (
        panel_box.left <= text_box.cx <= panel_box.right
        and panel_box.top <= text_box.cy <= panel_box.bottom
    )


def explicit_owner_key(name) -> Optional[str]:
    normalised = name.lower().replace(" ", "_")
    for suffix in (
        "_panel", "_card", "_footer", "_box", "_text", "_task", "_response",
        "_prompt", "_explanation", "_label",
    ):
        if normalised.endswith(suffix):
            return normalised[:-len(suffix)]
    return None


def likely_named_panel(name):
    normalised = name.lower().replace(" ", "_")
    return any(token in normalised for token in PANEL_NAME_TOKENS)


def likely_named_text(name):
    normalised = name.lower().replace(" ", "_")
    return any(token in normalised for token in TEXT_NAME_TOKENS)


def containment_metrics(text_box, panel_box, padding):
    distances = {
        "left": text_box.left - panel_box.left,
        "top": text_box.top - panel_box.top,
        "right": panel_box.right - text_box.right,
        "bottom": panel_box.bottom - text_box.bottom,
    }
    return {
        "inside": all(value >= 0 for value in distances.values()),
        "padded": all(value >= padding for value in distances.values()),
        "edge_clearance_in": {
            key: round(value / EMU_PER_INCH, 3)
            for key, value in distances.items()
        },
    }


def estimate_filled_panel_text_fit(shape):
    """Return a conservative rendered-height estimate for a filled text panel."""
    text_frame = shape.text_frame
    inner_width = shape.width - text_frame.margin_left - text_frame.margin_right
    inner_height = shape.height - text_frame.margin_top - text_frame.margin_bottom
    if inner_width <= 0 or inner_height <= 0:
        return None

    inner_width_pt = inner_width / 12700
    inner_height_pt = inner_height / 12700
    estimated_height_pt = 0.0
    measured_paragraphs = 0

    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        runs = [run for run in paragraph.runs if run.text]
        explicit_sizes = [run.font.size.pt for run in runs if run.font.size is not None]
        if not explicit_sizes:
            continue
        measured_paragraphs += 1
        default_size = max(explicit_sizes)
        estimated_width_pt = sum(
            len(run.text) * (run.font.size.pt if run.font.size is not None else default_size) * 0.5
            for run in runs
        )
        wrapped_lines = max(1, math.ceil(estimated_width_pt / inner_width_pt))
        estimated_height_pt += wrapped_lines * default_size * 1.2

    if measured_paragraphs == 0:
        return None
    return {
        "estimated_height_pt": round(estimated_height_pt, 1),
        "available_height_pt": round(inner_height_pt, 1),
        "estimated_fill_ratio": round(estimated_height_pt / inner_height_pt, 3),
    }
def choose_panel(text_shape, panels, text_index, indices):
    text_box = bbox(text_shape)
    text_name = name_of(text_shape)
    text_key = explicit_owner_key(text_name)
    candidates = []

    for panel in panels:
        if panel is text_shape or indices.get(id(panel), 10**9) >= text_index:
            continue
        panel_box = bbox(panel)
        ratio = overlap_ratio(text_box, panel_box)
        if ratio < MIN_PAIR_OVERLAP or not centre_inside(text_box, panel_box):
            continue

        panel_name = name_of(panel)
        panel_key = explicit_owner_key(panel_name)
        named_match = bool(text_key and panel_key and text_key == panel_key)
        name_signal = named_match or likely_named_panel(panel_name) or likely_named_text(text_name)
        score = (
            1 if named_match else 0,
            1 if name_signal else 0,
            ratio,
            -panel_box.area,
        )
        candidates.append((score, panel, ratio, named_match, name_signal))

    if not candidates:
        return None, {}

    candidates.sort(key=lambda item: item[0], reverse=True)
    _, panel, ratio, named_match, name_signal = candidates[0]
    confidence = (
        "explicit" if named_match
        else "strong" if ratio >= STRONG_PAIR_OVERLAP and name_signal
        else "spatial"
    )
    return panel, {
        "overlap_ratio": round(ratio, 4),
        "confidence": confidence,
        "named_match": named_match,
        "name_signal": name_signal,
    }


def audit_slide(slide, slide_no, slide_width, slide_height, padding_in):
    shapes = list(iter_shapes(slide.shapes))
    slide_area = slide_width * slide_height
    panels = [shape for shape in shapes if is_candidate_panel(shape, slide_area)]
    texts = [shape for shape in shapes if is_text_candidate(shape)]
    indices = {id(shape): index for index, shape in enumerate(shapes)}
    padding = int(round(padding_in * EMU_PER_INCH))
    pairs = []
    issues = []

    for text_shape in texts:
        if has_visible_fill(text_shape) and is_candidate_panel(text_shape, slide_area):
            text_frame = text_shape.text_frame
            margins = {
                "left": text_frame.margin_left,
                "right": text_frame.margin_right,
                "top": text_frame.margin_top,
                "bottom": text_frame.margin_bottom,
            }
            if any(value < padding for value in margins.values()):
                issues.append({
                    "severity": "fail",
                    "code": "internal_panel_margin_small",
                    "shape": name_of(text_shape),
                    "text": text_of(text_shape)[:160],
                    "message": f"Filled instructional text panel has internal margins below {padding_in:.2f} in.",
                    "margins_in": {
                        key: round(value / EMU_PER_INCH, 3)
                        for key, value in margins.items()
                    },
                })
            text_fit = estimate_filled_panel_text_fit(text_shape)
            if text_fit and text_fit["estimated_fill_ratio"] > 1.70:
                issues.append({
                    "severity": "fail",
                    "code": "internal_panel_text_fit_risk",
                    "slide": slide_no,
                    "shape": name_of(text_shape),
                    "text": text_of(text_shape)[:160],
                    "message": "Estimated wrapped text height substantially exceeds the filled panel's usable height.",
                    **text_fit,
                })

        panel, metadata = choose_panel(text_shape, panels, indices[id(text_shape)], indices)
        if panel is None:
            continue

        metrics = containment_metrics(bbox(text_shape), bbox(panel), padding)
        pair = {
            "text_shape": name_of(text_shape),
            "panel_shape": name_of(panel),
            "text": text_of(text_shape)[:160],
            **metadata,
            **metrics,
        }
        pairs.append(pair)

        if not metrics["inside"]:
            severity = "fail" if metadata["confidence"] in {"explicit", "strong"} else "warning"
            issues.append({
                "severity": severity,
                "code": "text_crosses_panel_bounds",
                **pair,
                "message": "Text box crosses the bounds of its likely containing panel.",
            })
        elif not metrics["padded"]:
            severity = "fail" if metadata["confidence"] == "explicit" else "warning"
            issues.append({
                "severity": severity,
                "code": "panel_padding_violation",
                **pair,
                "message": f"Text box is inside the panel but does not preserve {padding_in:.2f} in padding on every edge.",
            })

    status = (
        "fail" if any(issue["severity"] == "fail" for issue in issues)
        else "warning" if issues
        else "pass"
    )
    return {
        "slide": slide_no,
        "status": status,
        "candidate_panels": len(panels),
        "candidate_text_boxes": len(texts),
        "paired_text_boxes": len(pairs),
        "pairs": pairs,
        "issues": issues,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", required=True, help="PPTX file to audit")
    parser.add_argument("--out", required=True, help="JSON report path")
    parser.add_argument(
        "--padding-in",
        type=float,
        default=DEFAULT_PADDING_IN,
        help=f"Required nominal panel padding in inches (default {DEFAULT_PADDING_IN})",
    )
    args = parser.parse_args()

    presentation = Presentation(args.deck)
    results = [
        audit_slide(slide, index, presentation.slide_width, presentation.slide_height, args.padding_in)
        for index, slide in enumerate(presentation.slides, 1)
    ]
    candidate_panels = sum(result["candidate_panels"] for result in results)
    panel_slides = sum(result["candidate_panels"] > 0 for result in results)
    paired_text_boxes = sum(result["paired_text_boxes"] for result in results)
    coverage_issues = []
    if (
        candidate_panels >= MIN_PANEL_HEAVY_COUNT
        and panel_slides >= MIN_PANEL_HEAVY_SLIDES
        and paired_text_boxes == 0
    ):
        coverage_issues.append({
            "severity": "fail",
            "code": "ineffective_panel_audit_coverage",
            "message": (
                "The deck is panel-heavy but the audit paired zero text boxes "
                "with containing panels; containment was not meaningfully tested."
            ),
            "candidate_panels": candidate_panels,
            "panel_slides": panel_slides,
            "paired_text_boxes": paired_text_boxes,
        })

    summary = {
        "slides": len(results),
        "fail": sum(result["status"] == "fail" for result in results),
        "warning": sum(result["status"] == "warning" for result in results),
        "pass": sum(result["status"] == "pass" for result in results),
        "candidate_panels": candidate_panels,
        "panel_slides": panel_slides,
        "paired_text_boxes": paired_text_boxes,
        "coverage_failures": len(coverage_issues),
    }
    overall = (
        "fail" if summary["fail"] or coverage_issues
        else "warning" if summary["warning"]
        else "pass"
    )
    report = {
        "deck": args.deck,
        "artifact_sha256": file_sha256(Path(args.deck)),
        "overall_status": overall,
        "padding_in": args.padding_in,
        "summary": summary,
        "results": results,
        "coverage_issues": coverage_issues,
        "note": "Geometry screening only. Spatial pairings without explicit ownership are heuristic. Render-level inspection is mandatory.",
    }

    output = Path(args.out)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps({"overall_status": overall, **summary}))
    return 1 if overall == "fail" else 0


if __name__ == "__main__":
    raise SystemExit(main())
