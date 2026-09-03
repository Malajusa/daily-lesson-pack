#!/usr/bin/env python3
"""Audit projected readability and rough space use in a teaching PPTX.

This is deliberately a screening tool. It identifies likely typography and
space-utilisation problems for full-size human review; it does not claim to
measure pedagogical quality or exact visual occupancy.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import statistics
import re
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


WARMUP_MAIN_MIN_PT = 36.0
MEANINGFUL_HARD_FLOOR_PT = 24.0
STRUCTURAL_FLOOR_PT = 22.0
WHY_TARGET_FLOOR_PT = 28.0
LOW_OCCUPANCY = 0.35
VERY_LOW_OCCUPANCY = 0.25
GENERIC_EVIDENCE = re.compile(r"^(?:checked|reviewed|intentional|looks good|acceptable|pass)$", re.I)


def iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return sizes


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def text_of(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def is_incidental_label(text: str) -> bool:
    t = " ".join(text.lower().split())
    if t in {"all", "most", "some", "why", "answer", "model", "we do", "retrieval", "review"}:
        return True
    return len(t) <= 14 and len(t.split()) <= 2 and not any(ch in t for ch in "?=+−-×÷")


def is_title_shape(shape, slide_height: int) -> bool:
    # Treat a shallow text box in the top ~16% of the slide as likely title/header.
    return bool(getattr(shape, "has_text_frame", False) and shape.top < slide_height * 0.16 and shape.height < slide_height * 0.18)


def is_page_marker(shape, slide_height: int) -> bool:
    text = " ".join(text_of(shape).split())
    return bool(text.isdigit() and shape.top > slide_height * 0.82)


def warmup_detected(texts: list[str]) -> bool:
    joined = "\n".join(texts).lower()
    if "warm-up" in joined or "warm up" in joined:
        return True
    normalised = {" ".join(text.lower().split()).rstrip(":") for text in texts}
    markers = sum(1 for token in ("all", "most", "some", "why") if token in normalised)
    return markers >= 3


def area(shape) -> int:
    return max(0, int(shape.width)) * max(0, int(shape.height))


def audit_slide(slide, slide_index: int, slide_width: int, slide_height: int, forced_warmup: bool) -> dict:
    shapes = list(iter_shapes(slide.shapes))
    text_shapes = []
    texts = []

    for shape in shapes:
        text = text_of(shape)
        if text:
            texts.append(text)
            sizes = font_sizes(shape)
            text_shapes.append({
                "shape": shape,
                "text": text,
                "sizes": sizes,
                "median_pt": statistics.median(sizes) if sizes else None,
                "max_pt": max(sizes) if sizes else None,
                "title_like": is_title_shape(shape, slide_height),
            })

    warmup = forced_warmup or warmup_detected(texts)
    body_text = [
        x for x in text_shapes
        if not x["title_like"] and not is_page_marker(x["shape"], slide_height)
    ]
    known_body_sizes = [s for x in body_text for s in x["sizes"]]
    body_max = max(known_body_sizes) if known_body_sizes else None

    issues = []

    for item in body_text:
        text = " ".join(item["text"].split())
        if not text:
            continue

        structural = is_incidental_label(text)
        why_text = warmup and text.lower().startswith("why") and len(text.split()) > 3
        floor = STRUCTURAL_FLOOR_PT if structural else (
            WHY_TARGET_FLOOR_PT if why_text else (WARMUP_MAIN_MIN_PT if warmup else MEANINGFUL_HARD_FLOOR_PT)
        )
        if not item["sizes"]:
            issues.append({
                "severity": "fail",
                "code": "meaningful_font_size_unknown",
                "message": f"No explicit font size is available for a projected text box with a {floor:.0f} pt role floor: {text[:120]}",
            })
            continue
        smallest_pt = min(item["sizes"])
        if smallest_pt < floor:
            issues.append({
                "severity": "fail",
                "code": "undersized_body_text",
                "message": f"Smallest explicit run is {smallest_pt:.1f} pt (role floor {floor:.0f} pt): {text[:120]}",
            })

    # Rough occupancy: count visible non-title text, pictures, charts and tables.
    useful_area = float(slide_width * slide_height)
    occupied = 0
    for shape in shapes:
        if is_title_shape(shape, slide_height) or is_page_marker(shape, slide_height):
            continue
        if getattr(shape, "has_text_frame", False) and text_of(shape):
            occupied += area(shape)
        elif shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE}:
            occupied += area(shape)

    occupancy = min(1.0, occupied / useful_area) if useful_area else 0.0
    if occupancy < LOW_OCCUPANCY and (body_max is None or body_max <= 36.0):
        severity = "warning"
        if warmup and occupancy < VERY_LOW_OCCUPANCY and (body_max is None or body_max < 36.0):
            severity = "fail"
        issues.append({
            "severity": severity,
            "code": "possible_unused_space",
            "message": f"Estimated body occupancy is {occupancy:.0%} while body text remains relatively small. Inspect whether content could be enlarged or reflowed.",
        })

    status = "pass"
    if any(i["severity"] == "fail" for i in issues):
        status = "fail"
    elif issues:
        status = "warning"

    return {
        "slide": slide_index,
        "warmup_detected": warmup,
        "body_max_font_pt": body_max,
        "estimated_body_occupancy": round(occupancy, 4),
        "status": status,
        "issues": issues,
    }


def parse_slide_set(values: list[str]) -> set[int]:
    result: set[int] = set()
    for value in values:
        for part in value.split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                a, b = part.split("-", 1)
                result.update(range(int(a), int(b) + 1))
            else:
                result.add(int(part))
    return result


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", required=True, help="PPTX file to audit")
    parser.add_argument("--out", required=True, help="JSON report path")
    parser.add_argument(
        "--dispositions",
        help="Optional JSON ledger for individually reviewed warnings; it must match the deck SHA-256",
    )
    parser.add_argument(
        "--warmup-slides",
        action="append",
        default=[],
        help="Optional 1-based slide numbers/ranges to force as warm-ups, e.g. 2-11,35-54",
    )
    args = parser.parse_args()

    deck_path = Path(args.deck)
    out_path = Path(args.out)
    forced = parse_slide_set(args.warmup_slides)

    prs = Presentation(str(deck_path))
    results = [
        audit_slide(slide, idx, prs.slide_width, prs.slide_height, idx in forced)
        for idx, slide in enumerate(prs.slides, start=1)
    ]

    summary = {
        "slides": len(results),
        "fail": sum(r["status"] == "fail" for r in results),
        "warning": sum(r["status"] == "warning" for r in results),
        "pass": sum(r["status"] == "pass" for r in results),
    }
    deck_hash = sha256(deck_path)
    dispositions: dict[tuple[int, str], dict] = {}
    invalid_dispositions: list[str] = []
    if args.dispositions:
        try:
            ledger = json.loads(Path(args.dispositions).read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            ledger = {}
            invalid_dispositions.append(f"Disposition ledger is unreadable: {exc}")
        if str(ledger.get("artifact_sha256", "")).lower() != deck_hash:
            invalid_dispositions.append("Disposition ledger is not bound to this deck SHA-256.")
        for item in ledger.get("dispositions", []):
            if not isinstance(item, dict):
                invalid_dispositions.append("Every disposition must be an object.")
                continue
            evidence = str(item.get("evidence", "")).strip()
            required = all(str(item.get(field, "")).strip() for field in ("code", "decision", "reviewer", "run_id")) and len(evidence) >= 20 and not GENERIC_EVIDENCE.fullmatch(evidence)
            if not required or item.get("decision") != "accepted" or not isinstance(item.get("slide"), int):
                invalid_dispositions.append("A disposition needs slide, code, accepted decision, evidence, reviewer and run_id.")
                continue
            dispositions[(item["slide"], item["code"])] = item

    unresolved_warnings = 0
    resolved_warnings = 0
    for result in results:
        for item in result["issues"]:
            if item["severity"] != "warning":
                continue
            disposition = dispositions.get((result["slide"], item["code"]))
            if disposition:
                item["severity"] = "resolved_warning"
                item["disposition"] = disposition
                resolved_warnings += 1
            else:
                unresolved_warnings += 1

    hard_failures = sum(any(i["severity"] == "fail" for i in r["issues"]) for r in results)
    overall = "fail" if hard_failures or unresolved_warnings or invalid_dispositions else "pass"
    summary.update({
        "hard_failure_slides": hard_failures,
        "unresolved_warnings": unresolved_warnings,
        "resolved_warnings": resolved_warnings,
    })

    report = {
        "deck": str(deck_path),
        "overall_status": overall,
        "thresholds": {
            "warmup_main_min_pt": WARMUP_MAIN_MIN_PT,
            "meaningful_hard_floor_pt": MEANINGFUL_HARD_FLOOR_PT,
            "structural_floor_pt": STRUCTURAL_FLOOR_PT,
            "why_target_floor_pt": WHY_TARGET_FLOOR_PT,
            "low_occupancy": LOW_OCCUPANCY,
        },
        "summary": summary,
        "artifact_sha256": deck_hash,
        "dispositions": {
            "resolved_warnings": resolved_warnings,
            "unresolved_warnings": unresolved_warnings,
            "invalid": invalid_dispositions,
        },
        "results": results,
        "note": "Role-based screening. Every warning must be specifically dispositioned; full-size rendered inspection remains mandatory.",
    }

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps({"overall_status": overall, **summary}))
    return 1 if overall == "fail" else 0


if __name__ == "__main__":
    raise SystemExit(main())
