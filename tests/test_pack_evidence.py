import copy
import json
import sys
import tempfile
import unittest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import fitz

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0,str(ROOT/'scripts'))
from pack_evidence import audit_pack, audit_review, validate_content, digest, expected_checks
from audit_pack_contract import validate_sequence
from audit_slide_typography import audit_slide, font_sizes, is_incidental_label

class EvidenceTests(unittest.TestCase):
    def setUp(self):
        self.tmp=tempfile.TemporaryDirectory(); self.p=Path(self.tmp.name)
        self.content={'schema_version':3,'instances':[{'id':'lit','owner':'dlp-literacy-warmup','start':'08:30','duration_minutes':20,'purpose':'retrieval'}],
            'tasks':[{'id':'q1','instance_id':'lit','operation':'combine','fields':{'prompt':'Join the sentences using “but”: The day was hot. The night was cool.','answer':'The day was hot, but the night was cool.'},'demands':[{'id':'response','action':'combine','answer_quote':'The day was hot, but the night was cool.'}]}],
            'documents':{'brief':{'fields':{'text':'Teach the contrast conjunction “but”.'}}}}
        self.context={'timetable_instances':self.content['instances'],'required_artifacts':['deck','briefing','student','answers']}
        self.manifest={'schema_version':3,'artifacts':[],'bindings':[]}
        prs=Presentation()
        for field in ('prompt','answer'):
            slide=prs.slides.add_slide(prs.slide_layouts[6]); shape=slide.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(3));shape.text=self.content['tasks'][0]['fields'][field]
            self.manifest['bindings'].append({'artifact':'deck','record':'q1','field':field,'page':len(prs.slides),'shape_id':shape.shape_id})
        prs.save(self.p/'deck.pptx')
        self.manifest['artifacts'].append({'id':'deck','role':'deck','path':'deck.pptx','sha256':digest(self.p/'deck.pptx'),'pages':2})
        for role,record,field in [('briefing','brief','text'),('student','q1','prompt'),('answers','q1','answer')]:
            text=self.content['documents']['brief']['fields']['text'] if record=='brief' else self.content['tasks'][0]['fields'][field]
            doc=fitz.open(); page=doc.new_page();page.insert_textbox(fitz.Rect(40,40,550,200),text,fontsize=12,fontname='helv')
            # Standard PDF font cannot encode curly quotes: use the exact exported text as canonical fixture.
            text=page.get_text(clip=fitz.Rect(40,40,550,200)).strip()
            if record=='brief': self.content['documents']['brief']['fields'][field]=text
            else:
                self.content['tasks'][0]['fields'][field]=text
            doc.save(self.p/(role+'.pdf'));doc.close()
            self.manifest['artifacts'].append({'id':role,'role':role,'path':role+'.pdf','sha256':digest(self.p/(role+'.pdf')),'pages':1})
            self.manifest['bindings'].append({'artifact':role,'record':record,'field':field,'page':1,'bbox':[40,40,550,200]})
        # Rebuild deck from final canonical fixture values.
        for slide,field in zip(prs.slides,('prompt','answer')): slide.shapes[0].text=self.content['tasks'][0]['fields'][field]
        prs.save(self.p/'deck.pptx');self.manifest['artifacts'][0]['sha256']=digest(self.p/'deck.pptx')
        self.save()
    def tearDown(self): self.tmp.cleanup()
    def save(self):
        (self.p/'content.json').write_text(json.dumps(self.content));(self.p/'context.json').write_text(json.dumps(self.context))
        self.manifest['content_sha256']=digest(self.p/'content.json');self.manifest['context_sha256']=digest(self.p/'context.json')
        (self.p/'manifest.json').write_text(json.dumps(self.manifest))
    def audit(self): return audit_pack(self.p/'manifest.json',self.p/'content.json',self.p/'context.json')[0]
    def test_valid_complete_pack(self): self.assertEqual(self.audit(),[])
    def test_changed_printable(self):
        doc=fitz.open();page=doc.new_page();page.insert_text((40,60),'Write A or B and name the focus.');doc.save(self.p/'student.pdf');doc.close()
        self.manifest['artifacts'][2]['sha256']=digest(self.p/'student.pdf');self.save()
        self.assertTrue(any('Content mismatch' in e for e in self.audit()))
    def test_unbound_printable_instruction(self):
        doc=fitz.open(self.p/'student.pdf');doc[0].insert_text((40,300),'Also explain your answer and name the conjunction.');doc.save(self.p/'changed.pdf');doc.close()
        (self.p/'changed.pdf').replace(self.p/'student.pdf');self.manifest['artifacts'][2]['sha256']=digest(self.p/'student.pdf');self.save()
        self.assertTrue(any('Unbound PDF text' in e for e in self.audit()))
    def test_main_role_above_body(self):
        prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);shape=slide.shapes.add_textbox(Inches(1),Inches(.2),Inches(7),Inches(.6));shape.text='Write the completed sentence.';shape.name='DLP:main';shape.text_frame.paragraphs[0].runs[0].font.size=Pt(8)
        r=audit_slide(slide,1,prs.slide_width,prs.slide_height,True)
        self.assertTrue(any(i['code']=='undersized_body_text' for i in r['issues']))
    def test_bottom_numeric_answer_keeps_role_floor(self):
        prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);shape=slide.shapes.add_textbox(Inches(1),int(prs.slide_height*.85),Inches(2),Inches(.3));shape.text='42';shape.name='DLP:main';shape.text_frame.paragraphs[0].runs[0].font.size=Pt(8)
        r=audit_slide(slide,1,prs.slide_width,prs.slide_height,True)
        self.assertTrue(any(i['code']=='undersized_body_text' for i in r['issues']))
    def test_actual_response_override(self):
        from audit_pack_contract import audit_deck
        prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);box=slide.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(3));box.text='LITERACY WARM-UP 1 OF 1 QUESTION\nJoin the sentences and explain why.'
        prs.save(self.p/'override.pptx')
        issues,_=audit_deck(str(self.p/'override.pptx'),1,0,{1})
        self.assertFalse(any(i['code']=='literacy_extra_reasoning_demand' for i in issues))
    def test_stale_answer_pdf(self):
        with (self.p/'answers.pdf').open('ab') as f:f.write(b'\nchanged')
        self.assertTrue(any('stale hash' in e for e in self.audit()))
    def test_missing_required_file(self):
        self.manifest['artifacts'].pop();self.manifest['bindings'].pop();self.save()
        self.assertTrue(any('deliverable' in e for e in self.audit()))
    def test_stale_content(self):
        (self.p/'content.json').write_text(json.dumps(dict(self.content,note='changed')))
        self.assertIn('Stale content/context binding',self.audit())
    def test_duplicate_id(self):
        self.content['tasks'].append(copy.deepcopy(self.content['tasks'][0]))
        self.assertIn('Duplicate instance/task ID',validate_content(self.content))
    def test_extra_response(self):
        t=self.content['tasks'][0];t['fields']['prompt']='Write A or B and name the focus.'
        self.assertTrue(any('additional response' in e for e in validate_content(self.content)))
    def test_generic_combine_with_but(self):
        self.content['tasks'][0]['fields']['prompt']='Write a sentence using “but”.'
        self.assertTrue(any('exact action' in e for e in validate_content(self.content)))
    def test_implausible_distractor(self):
        t=self.content['tasks'][0];t['fields']['prompt']='Choose Across these habitats or For no reason.'
        t['options']=[{'text':'Across these habitats','correct':True},{'text':'For no reason','correct':False,'misconception':'invented','rejection_reason':'invented'}]
        self.assertTrue(any('implausible' in e for e in validate_content(self.content)))
    def test_missing_proposition(self):
        t=self.content['tasks'][0];t['operation']='revise';t['fields'].update(before='Wetlands are valuable habitats.',after='These habitats reduce flooding.')
        t['propositions']=[{'before_quote':'valuable habitats','after_quote':'valuable habitats'}]
        self.assertTrue(any('proposition lost' in e for e in validate_content(self.content)))
    def test_valid_proposition(self):
        t=self.content['tasks'][0];t['operation']='revise';t['fields'].update(before='Wetlands are valuable habitats.',after='These wetlands are valuable habitats.')
        t['propositions']=[{'before_quote':'valuable habitats','after_quote':'valuable habitats'}]
        self.assertEqual(validate_content(self.content),[])
    def test_absent_unscheduled_warmup(self): self.assertEqual(validate_sequence([],['QUESTION','ANSWER'],'numeracy',0),[])
    def test_missing_scheduled_warmup(self): self.assertTrue(validate_sequence([],['QUESTION','ANSWER'],'numeracy',5))
    def test_extra_unscheduled_warmup(self):
        self.assertTrue(validate_sequence([{'number':1,'total':1,'role':'QUESTION','slide':1}],['QUESTION','ANSWER'],'numeracy',0))
    def make_review(self):
        binding=self.manifest['bindings'][0]
        citation={k:v for k,v in binding.items() if k in ('artifact','page','shape_id')};citation['quote']=self.content['tasks'][0]['fields']['prompt']
        checks=[{'id':i,'subject':s,'result':'PASS','observation':'Synthetic fixture observation','citations':[citation]} for i,s in expected_checks(self.content,self.manifest,'semantic')]
        review={'manifest_sha256':digest(self.p/'manifest.json'),'content_sha256':digest(self.p/'content.json'),'requirements_sha256':digest(ROOT/'references/qa-requirements.json'),'execution_id':'execution-2','checks':checks}
        (self.p/'review.json').write_text(json.dumps(review));(self.p/'host.txt').write_text('Synthetic host execution fixture')
        trace={'source':'external-runner','execution_id':'execution-2','reviewer_actor':'reviewer','generator_actor':'generator','review_sha256':digest(self.p/'review.json'),'manifest_sha256':digest(self.p/'manifest.json'),'transcript_path':'host.txt','transcript_sha256':digest(self.p/'host.txt')}
        (self.p/'trace.json').write_text(json.dumps(trace));return review,trace
    def review_audit(self):return audit_review(self.p/'review.json','semantic',self.p/'manifest.json',self.p/'content.json',self.content,self.manifest,self.p/'trace.json')
    def test_release_command_runs_audits_and_accepts_v3_na(self):
        from unittest.mock import patch
        from contextlib import redirect_stdout
        import io
        import audit_release_bundle as gate
        # This integration test isolates the already separately-tested structural
        # audit executables; pack parsing, file comparison and both reviews are real.
        self.manifest['renders']=[]
        for artifact in self.manifest['artifacts']:
            for n in range(1,artifact['pages']+1):
                path=self.p/(artifact['id']+str(n)+'.png')
                pix=fitz.Pixmap(fitz.csRGB,fitz.IRect(0,0,10,10),False);pix.clear_with(255);pix.save(path)
                self.manifest['renders'].append({'artifact':artifact['id'],'page':n,'path':path.name,'sha256':digest(path),'artifact_sha256':artifact['sha256']})
        self.save();review,trace=self.make_review()
        check=next(c for c in review['checks'] if c['id']=='TASK.DISTRACTORS')
        check.update(result='NA',applicability_reason='Direct production task has no options.')
        (self.p/'review.json').write_text(json.dumps(review));trace['review_sha256']=digest(self.p/'review.json');(self.p/'trace.json').write_text(json.dumps(trace))
        visual=copy.deepcopy(review);visual['checks']=[{'id':i,'subject':sub,'result':'PASS','observation':'Synthetic rendered page fixture','citations':[{'artifact':sub.rsplit(':',1)[0],'page':int(sub.rsplit(':',1)[1])}]} for i,sub in expected_checks(self.content,self.manifest,'visual')]
        (self.p/'visual-review.json').write_text(json.dumps(visual));vtrace=dict(trace,review_sha256=digest(self.p/'visual-review.json'));(self.p/'visual-trace.json').write_text(json.dumps(vtrace))
        (self.p/'components.json').write_text(json.dumps({'generation_run_id':'generation'}));(self.p/'warnings.json').write_text('{}')
        args=['gate','--deck',str(self.p/'deck.pptx'),'--out',str(self.p/'release.json')]
        for flag in ['contract','year-profile','typography','containment','visual']:
            args+=['--'+flag,str(self.p/(flag+'.json'))]
        for flag,name in [('manifest','manifest.json'),('content','content.json'),('context-record','context.json'),('component-record','components.json'),('warning-ledger','warnings.json'),('visual-review','visual-review.json'),('semantic-review','review.json'),('semantic-trace','trace.json'),('visual-trace','visual-trace.json')]:args+=['--'+flag,str(self.p/name)]
        def screen(command,**kwargs):
            Path(command[command.index('--out')+1]).write_text(json.dumps({'status':'PASS','artifact_sha256':digest(self.p/'deck.pptx')}))
            return type('Result',(),{'returncode':0,'stderr':''})()
        with patch.object(sys,'argv',args),patch.object(gate.subprocess,'run',side_effect=screen) as screens,redirect_stdout(io.StringIO()):
            self.assertEqual(gate.main(),0)
            self.assertEqual(screens.call_count,5)
        self.assertEqual(json.loads((self.p/'release.json').read_text())['status'],'PASS')
        # The same supplied PASS files cannot rescue an incomplete new review.
        review['checks']=[];(self.p/'review.json').write_text(json.dumps(review));trace['review_sha256']=digest(self.p/'review.json');(self.p/'trace.json').write_text(json.dumps(trace))
        with patch.object(sys,'argv',args),patch.object(gate.subprocess,'run') as screens,redirect_stdout(io.StringIO()):
            self.assertEqual(gate.main(),1)
            self.assertEqual(screens.call_count,0)
        self.assertEqual(json.loads((self.p/'release.json').read_text())['status'],'FAIL')
    def test_complete_review_structure(self): self.make_review();self.assertEqual(self.review_audit(),[])
    def test_unscoped_semantic_assertion(self):
        review,trace=self.make_review();review['checks']=[]
        (self.p/'review.json').write_text(json.dumps(review));trace['review_sha256']=digest(self.p/'review.json');(self.p/'trace.json').write_text(json.dumps(trace))
        self.assertTrue(any('missing check' in e for e in self.review_audit()))
    def test_same_actor_fails(self):
        _,trace=self.make_review();trace['reviewer_actor']='generator';(self.p/'trace.json').write_text(json.dumps(trace))
        self.assertTrue(any('separate reviewer' in e for e in self.review_audit()))
    def test_stale_review(self):
        self.make_review();self.manifest['note']='changed';self.save()
        self.assertTrue(any('different pack' in e for e in self.review_audit()))
    def test_missing_pdf_visual_coverage(self):
        expected=expected_checks(self.content,self.manifest,'visual')
        self.assertIn(('VISUAL.PAGE','answers:1'),expected)
    def test_effective_paragraph_font(self):
        prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);shape=slide.shapes.add_textbox(Inches(1),Inches(2),Inches(7),Inches(2));shape.text='A meaningful prompt'
        shape.text_frame.paragraphs[0].font.size=Pt(36)
        self.assertEqual(font_sizes(shape),[36])
    def test_subordinate_floor(self):
        prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);shape=slide.shapes.add_textbox(Inches(1),Inches(2),Inches(7),Inches(3));shape.text='Write the completed sentence.';shape.name='DLP:instruction';shape.text_frame.paragraphs[0].runs[0].font.size=Pt(28)
        r=audit_slide(slide,1,prs.slide_width,prs.slide_height,True)
        self.assertFalse(any(i['code']=='undersized_body_text' for i in r['issues']))
        shape.name='DLP:main';r=audit_slide(slide,1,prs.slide_width,prs.slide_height,True)
        self.assertTrue(any(i['code']=='undersized_body_text' for i in r['issues']))
    def test_decimal_answer_not_structural(self):self.assertFalse(is_incidental_label('0.75'))

if __name__=='__main__':unittest.main()
