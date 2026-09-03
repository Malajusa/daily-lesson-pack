from __future__ import annotations

import hashlib
import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


def load_module(name: str, relative: str):
    path = ROOT / relative
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def add_text(slide, text: str, top: float, size: float = 36) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(11.5), Inches(0.7))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)


class HardeningContractTests(unittest.TestCase):
    def test_default_profile_requires_five_numeracy_pairs(self) -> None:
        profile = json.loads((ROOT / "references/default-pack-profile.json").read_text())
        self.assertEqual(profile["numeracy_warmup"]["prompt_answer_pairs"], 5)
        self.assertIn("10 slides total", (ROOT / "skills/dlp-numeracy-warmup/SKILL.md").read_text())

    def test_ten_pair_numeracy_deck_is_rejected(self) -> None:
        audit = load_module("contract_audit_numeracy", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            path = Path(folder) / "ten-pairs.pptx"
            prs = Presentation()
            prs.slides._sldIdLst.clear()
            for number in range(1, 11):
                for role in ("QUESTION", "ANSWER"):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_text(slide, f"NUMERACY WARM-UP {number} OF 10 — {role}", 0.2)
            prs.save(path)
            issues, _ = audit.audit_deck(str(path), literacy_count=1, numeracy_count=5)
        codes = {item["code"] for item in issues}
        self.assertIn("numeracy_warmup_count_invalid", codes)

    def test_repeated_maths_owner_passes_with_unique_instances_and_evidence(self) -> None:
        audit = load_module("contract_audit_instances", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            folder_path = Path(folder)
            deck = folder_path / "deck.pptx"
            prs = Presentation()
            prs.save(deck)
            digest = hashlib.sha256(deck.read_bytes()).hexdigest()
            instances = [
                {"id": "maths-1", "owner": "dlp-maths-lesson", "start": "08:30", "duration_minutes": 60, "purpose": "explicit teaching"},
                {"id": "maths-2", "owner": "dlp-maths-lesson", "start": "10:50", "duration_minutes": 60, "purpose": "application and exit"},
            ]
            context = folder_path / "context.json"
            context.write_text(json.dumps({
                "active_year_profile": {"value": "year-4-5"},
                "mathematics_focus": {"value": "whole-number place value"},
                "timetable_instances": instances,
            }))
            required = set(audit.REQUIRED_CHECKS["dlp-maths-lesson"])
            required.update({"MATHS.YEAR4.PATHWAY", "MATHS.YEAR5.PATHWAY"})
            required.add("MATHS.BLOCK.BREAKPOINT")
            record = {
                "schema_version": 2,
                "generation_run_id": "generation-test-1",
                "artifact_sha256": digest,
                "scheduled_instances": instances,
                "components": [
                    {
                        "instance_id": item["id"],
                        "owner": item["owner"],
                        "active_year_profile": "year-4-5",
                        "status": "PASS",
                        "estimated_minutes": 55,
                        "slide_range": "1-8",
                        "checks": [{"id": check, "result": "PASS", "evidence": "slides 1-8 inspected"} for check in sorted(required)],
                    }
                    for item in instances
                ],
            }
            record_path = folder_path / "components.json"
            record_path.write_text(json.dumps(record))
            issues, summary = audit.audit_component_record(str(record_path), str(context), str(deck))
        self.assertEqual(issues, [])
        self.assertEqual(summary["passed"], 2)

    def test_freeform_component_pass_is_rejected(self) -> None:
        audit = load_module("contract_audit_self_attestation", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            path = Path(folder) / "record.json"
            path.write_text(json.dumps({"scheduled_components": ["dlp-maths-lesson"], "components": [{"name": "dlp-maths-lesson", "status": "PASS", "checks": ["checked"], "artefact": "slides 1-9"}]}))
            issues, _ = audit.audit_component_record(str(path))
        codes = {item["code"] for item in issues}
        self.assertIn("component_record_schema_outdated", codes)
        self.assertIn("scheduled_instances_missing", codes)

    def test_component_time_overrun_is_rejected(self) -> None:
        audit = load_module("contract_audit_time", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            folder_path = Path(folder)
            deck = folder_path / "deck.pptx"
            Presentation().save(deck)
            digest = hashlib.sha256(deck.read_bytes()).hexdigest()
            instance = {"id": "numeracy-1", "owner": "dlp-numeracy-warmup", "start": "08:40", "duration_minutes": 10, "purpose": "retrieval"}
            required = audit.REQUIRED_CHECKS["dlp-numeracy-warmup"]
            record = {
                "schema_version": 2,
                "generation_run_id": "generation-test-2",
                "artifact_sha256": digest,
                "scheduled_instances": [instance],
                "components": [{
                    "instance_id": "numeracy-1",
                    "owner": "dlp-numeracy-warmup",
                    "status": "PASS",
                    "estimated_minutes": 18,
                    "slide_range": "1-10",
                    "checks": [{"id": item, "result": "PASS", "evidence": "slides 1-10"} for item in sorted(required)],
                }],
            }
            record_path = folder_path / "record.json"
            record_path.write_text(json.dumps(record))
            issues, _ = audit.audit_component_record(str(record_path), deck_path=str(deck))
        self.assertIn("component_time_budget_exceeded", {item["code"] for item in issues})

    def test_year_profile_audit_accepts_instance_schema(self) -> None:
        with tempfile.TemporaryDirectory() as folder:
            folder_path = Path(folder)
            deck = folder_path / "deck.pptx"
            Presentation().save(deck)
            digest = hashlib.sha256(deck.read_bytes()).hexdigest()
            context = folder_path / "context.json"
            context.write_text(json.dumps({
                "active_year_profile": {
                    "value": "year-4-5",
                    "source": "references/year-level-profiles/year-4-5.md",
                    "resolved": True,
                    "status": "calibrated",
                }
            }))
            component = folder_path / "component.json"
            component.write_text(json.dumps({
                "artifact_sha256": digest,
                "scheduled_instances": [{"id": "maths-1", "owner": "dlp-maths-lesson"}],
                "components": [{
                    "instance_id": "maths-1",
                    "owner": "dlp-maths-lesson",
                    "active_year_profile": "year-4-5",
                }],
            }))
            result = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts/audit_year_profile_context.py"),
                    "--context-record", str(context),
                    "--component-record", str(component),
                    "--deck", str(deck),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
        self.assertEqual(result.returncode, 0, result.stdout + result.stderr)
        self.assertEqual(json.loads(result.stdout)["status"], "PASS")

    def test_year6_does_not_require_year45_pathway_checks(self) -> None:
        audit = load_module("contract_audit_year6", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            folder_path = Path(folder)
            deck = folder_path / "deck.pptx"
            Presentation().save(deck)
            digest = hashlib.sha256(deck.read_bytes()).hexdigest()
            instance = {"id": "maths-1", "owner": "dlp-maths-lesson", "start": "09:00", "duration_minutes": 60, "purpose": "explicit teaching"}
            context = folder_path / "context.json"
            context.write_text(json.dumps({
                "active_year_profile": {"value": "year-6"},
                "mathematics_focus": {"value": "multiplicative thinking"},
                "timetable_instances": [instance],
            }))
            required = set(audit.REQUIRED_CHECKS["dlp-maths-lesson"])
            record = folder_path / "record.json"
            record.write_text(json.dumps({
                "schema_version": 2,
                "generation_run_id": "generation-year6",
                "artifact_sha256": digest,
                "scheduled_instances": [instance],
                "components": [{
                    "instance_id": "maths-1",
                    "owner": "dlp-maths-lesson",
                    "active_year_profile": "year-6",
                    "status": "PASS",
                    "estimated_minutes": 55,
                    "slide_range": "1-8",
                    "checks": [{"id": check, "result": "PASS", "evidence": "slides 1-8 inspected"} for check in sorted(required)],
                }],
            }))
            issues, _ = audit.audit_component_record(str(record), str(context), str(deck))
        missing = {item.get("check_id") for item in issues if item["code"] == "component_required_check_missing"}
        self.assertNotIn("MATHS.YEAR4.PATHWAY", missing)
        self.assertNotIn("MATHS.YEAR5.PATHWAY", missing)

    def test_small_warmup_element_fails_even_when_another_is_36pt(self) -> None:
        audit = load_module("typography_audit_roles", "scripts/audit_slide_typography.py")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text(slide, "NUMERACY WARM-UP 1 OF 5 — QUESTION", 0.1, 36)
        add_text(slide, "All", 1.0, 24)
        add_text(slide, "Calculate 245 + 37", 2.0, 36)
        add_text(slide, "Explain which place value changed.", 3.0, 24)
        result = audit.audit_slide(slide, 1, prs.slide_width, prs.slide_height, True)
        failures = [item for item in result["issues"] if item["code"] == "undersized_body_text"]
        self.assertTrue(failures)
        self.assertEqual(result["status"], "fail")

    def test_fraction_slash_notation_is_rejected_on_maths_slide(self) -> None:
        audit = load_module("contract_audit_fractions", "scripts/audit_pack_contract.py")
        with tempfile.TemporaryDirectory() as folder:
            path = Path(folder) / "slash.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_text(slide, "MATHEMATICS — Show why 3/5 equals 60/100", 1.0)
            prs.save(path)
            issues, _ = audit.audit_deck(str(path), literacy_count=1, numeracy_count=1)
        self.assertIn("projected_fraction_slash_notation", {item["code"] for item in issues})

    def test_unresolved_typography_warning_blocks_release(self) -> None:
        with tempfile.TemporaryDirectory() as folder:
            folder_path = Path(folder)
            deck = folder_path / "warning.pptx"
            report = folder_path / "report.json"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_text(slide, "NUMERACY WARM-UP 1 OF 5 — QUESTION", 0.1, 36)
            add_text(slide, "Calculate 5 + 7", 2.0, 36)
            prs.save(deck)
            result = subprocess.run(
                [sys.executable, str(ROOT / "scripts/audit_slide_typography.py"), "--deck", str(deck), "--out", str(report)],
                capture_output=True,
                text=True,
                check=False,
            )
            payload = json.loads(report.read_text())
        self.assertNotEqual(result.returncode, 0)
        self.assertGreater(payload["dispositions"]["unresolved_warnings"], 0)


if __name__ == "__main__":
    unittest.main()
